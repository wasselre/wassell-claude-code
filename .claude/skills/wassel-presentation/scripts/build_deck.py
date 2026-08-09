"""
Wassel Real Estate — deterministic presentation builder.

Produces a 16-slide 16:9 Arabic RTL PowerPoint in Wassel's brand style.
The layout is fixed. Content comes in via a dict (see SCHEMA below).

Usage:
    from build_deck import build
    build(content, output_path="/mnt/user-data/outputs/my_deck.pptx")

SCHEMA (every field is a string unless noted):
    content = {
        # --- identifiers (used across slides) ---
        "project_name":       "مقام كورتيارد 17",
        "district":           "حي النرجس",
        "city":               "الرياض",
        "year":               "2026",
        "footer_right":       "وصل العقارية  |  مقام كورتيارد 17",   # appears on every content slide
        "marketing_footer_right": "وصل العقارية  |  الخطة التسويقية",  # slides 8-10
        "sales_footer_right":     "وصل العقارية  |  الخطة البيعية",    # slides 12-15

        # --- slide 4 (market analysis) ---
        # KPI tiles are CONTENT-DRIVEN. Pick the three most presentable,
        # unambiguously-scoped metrics from the research — don't force a
        # specific indicator. No hardcoded slot labels like "إجمالي الوحدات في المربع"
        # or "متوسط نسبة البيع" — if that number isn't in the data, pick a
        # different metric instead of leaving a dead "—" tile.
        # Each label must state its scope (e.g. "... في الحي" for district-level,
        # "... للمشروع" for project-level) so no tile can be misread.
        # TEXT RULE: no parentheses `()` anywhere in market.insight_lines or
        # any other copy — use em-dashes `— text —` around clarifications.
        "market": {
            "kpis": [                                       # EXACTLY 3 tiles, each: {"value": "...", "label": "..."}
                {"value": "10,202", "label": "متوسط سعر متر الدور في الحي"},
                {"value": "173",    "label": "متوسط مساحة الدور في الحي"},
                {"value": "7",      "label": "صفقات الدور خلال 12 شهر"},
            ],
            "insight_lines":    ["...", "..."],          # EXACTLY 2 lines, each ≤ 12 words (ideally ≤ 8). Card is 9.10"×1.10" at 14pt — full sentences will overflow. Condense to punchy one-liners. No parens.
            "price_range_min":  "830,000",
            "price_range_max":  "1,500,000",
            "avg_price":        "1,200,000",
        },

        # --- slide 5 (competitors table) ---
        "competitors": {
            "headers":  ["المشروع", "المطوّر", "عدد الوحدات", "نسبة البيع", "متوسط السعر", "ملاحظات"],
            "rows":     [ [...6 cells...], ... ],
            "footnote": "* ...",
        },

        # --- slide 6 (project + opportunity) ---
        "project": {
            "location":         "حي النرجس —  دقيقتين من طريق الملك سلمان",
            "area_range":       "91 م² حتى 265 م²",
            "price_range":      "918,400 الى 1,554,170 ريال",
            "unit_total":       "61 وحدة — 4 عمائر \u200E(A/B/C/D)\u200E",  # parens wrapped with LRM
            "units_available":  "~43 وحدة",
            "amenities":        "نادي رياضي، لاونج، حضانة، غرفة سائق، مواقف",
            "warranties":       "هيكل إنشائي 10 سنوات  |    عزل 15 سنة",
            "opportunity_bullets": [
                "الطلب مرتفع: الحي يسجل 7 وحدات مباعة شهرياً",
                "المرافق تُحرّك البيع — مقام 17 يملكها",
                "معدل مبيعات المشروع قريب من متوسط السوق",
                "43 وحدة متاحة بقيمة تجاوز 50 مليون ريال",
            ],
            "expected_revenue":   "+50,000,000",
            "remaining_units":    "43",
            "avg_price_millions": "1.2",
        },

        # --- slide 10 (marketing targets) ---
        # ONLY total_units_sold is required. The view-target and lead-target
        # numbers displayed are DERIVED from this via FIXED conversion rates
        # (view→lead = 1%, lead→sale = 0.6%). Do not pass per-project
        # overrides — the derived-funnel narrative is a brand rule.
        "marketing_targets": {
            "total_units_sold":  "43",
        },

        # --- slide 14 (sales plan by the numbers, formula-driven table) ---
        # PRIMARY INPUT for the whole slide. Every row in the table derives
        # from this number via fixed constants (6%, 40%, 20%, 60%, 80%,
        # natural_visits = 2 × appointment visits). Round to nearest integer.
        "sales_plan": {
            "leads_per_month":  "1500",
        },

        # --- slide 14 right-sidebar KPIs (content-driven) ---
        # These are displayed as-is, separate from the funnel table.
        "funnel_kpis": {
            "overall_conversion":  "0.61%",
            "leads_per_sale":      "166",
            "avg_unit_price":      "1.3M",
            "months_to_sellout":   "3 - 6",
        },
    }

Missing keys are tolerated — the helper functions substitute a placeholder like "—"
and the slide structure stays intact. Use that rather than inventing numbers.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ─── Brand constants (do not change) ──────────────────────────────────────────
COPPER   = RGBColor(0xB8, 0x73, 0x4F)
SAND     = RGBColor(0xE8, 0xD9, 0xC0)
BROWN    = RGBColor(0x6B, 0x42, 0x26)
CREAM    = RGBColor(0xF8, 0xF5, 0xE9)
GOLD     = RGBColor(0xD9, 0xB5, 0x7F)
CHARCOAL = RGBColor(0x3F, 0x3F, 0x3F)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Amiri"

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(HERE, "..", "assets", "wassel_logo_white.png")

# ─── Low-level helpers ────────────────────────────────────────────────────────

def _add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
    """Plain filled rectangle. No outline by default."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
    shp.shadow.inherit = False
    return shp

def _set_rtl(paragraph):
    """Set paragraph to right-to-left. python-pptx doesn't expose this directly."""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")

def _set_font_all_scripts(run, font_name):
    """Set Latin, East-Asian, AND Complex-Script font on a run.

    WHY: `run.font.name = "Amiri"` only writes <a:latin> in the XML. PowerPoint
    uses the Complex-Script slot (<a:cs>) when rendering Arabic — if that slot
    is empty, Arabic silently falls back to the default theme font and the dropdown
    lies about what's actually rendering. Setting all three slots makes Amiri
    stick for every character class without the user having to re-click the font.
    """
    run.font.name = font_name  # writes <a:latin>
    rPr = run._r.get_or_add_rPr()
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    # Remove any stale ea/cs so we own them
    for local in ("ea", "cs"):
        for existing in rPr.findall(qn(f"a:{local}")):
            rPr.remove(existing)
    # Insert ea then cs after latin (schema order: latin, ea, cs, sym...)
    latin = rPr.find(qn("a:latin"))
    insert_at = list(rPr).index(latin) + 1 if latin is not None else len(rPr)
    for local in ("ea", "cs"):
        elem = etree.Element(qn(f"a:{local}"))
        elem.set("typeface", font_name)
        rPr.insert(insert_at, elem)
        insert_at += 1

# Unicode bidi marks
LRM = "\u200E"   # Left-to-right mark — locks LTR for one neutral character
RLM = "\u200F"   # Right-to-left mark — locks RTL for one neutral character
ALM = "\u061C"   # Arabic letter mark — strong RTL, zero-width
LRE = "\u202A"   # Left-to-right EMBEDDING — opens an LTR scope
PDF = "\u202C"   # Pop directional format — closes the embedding scope
NBSP = "\u00A0"  # Non-breaking space

import re as _re

def _normalize_spacing(text):
    """Clean up spacing around separators so the area before and after is equal.

    Rule: for any separator (—, –, -, _, |) that already has whitespace on at
    least one side (meaning it's being used as a separator, not as a compound
    joiner like 'A-B' or 'snake_case'), enforce exactly one space before and
    one space after. In Arabic-containing text, also wrap the separator in RLM
    marks so PowerPoint's bidi algorithm doesn't collapse the surrounding
    spaces when resolving the bidi-neutral separator against Arabic letters
    and Arabic-Indic digits."""
    if not isinstance(text, str) or not text:
        return text
    text = _re.sub(r" {2,}", " ", text)
    has_arabic = any("\u0600" <= ch <= "\u06FF" for ch in text)
    # Em-dash and en-dash are always separators (never appear inside compound
    # tokens), so match greedily with any surrounding whitespace (including none).
    for sep in ("—", "–"):
        pat_sep = _re.escape(sep)
        if has_arabic:
            text = _re.sub(rf"\s*{pat_sep}\s*", f" {RLM}{sep}{RLM} ", text)
        else:
            text = _re.sub(rf"\s*{pat_sep}\s*", f" {sep} ", text)
    # Hyphen-minus and underscore: only normalize when they already have
    # whitespace on at least one side. If neither side has whitespace, the
    # separator is part of a compound token (e.g. "A-B", "snake_case") and
    # must be left alone.
    for sep in ("-", "_"):
        pat_sep = _re.escape(sep)
        pattern = rf"(?:\s+{pat_sep}\s*|\s*{pat_sep}\s+)"
        if has_arabic:
            text = _re.sub(pattern, f" {RLM}{sep}{RLM} ", text)
        else:
            text = _re.sub(pattern, f" {sep} ", text)
    text = _re.sub(r"\s*\|\s*", " | ", text)
    text = _re.sub(r" {2,}", " ", text)
    return text.strip()

# Western → Arabic-Indic digit translation table.
# U+0660..U+0669 = ٠ ١ ٢ ٣ ٤ ٥ ٦ ٧ ٨ ٩
_WESTERN_TO_ARABIC_INDIC = str.maketrans("0123456789", "٠١٢٣٤٥٦٧٨٩")

def _convert_digits(text):
    """Convert Western digits (0-9) to Arabic-Indic digits (٠-٩).

    Also replace decimal-point dots (between digits) with the Arabic decimal
    separator ٫ (U+066B). WHY: a Latin "." is a bidi-neutral character, so
    inside an RTL Arabic paragraph a decimal number like "0.6%" can have its
    "." drift to the wrong visual side of the digits. The Arabic decimal
    separator ٫ is a strong-RTL character that stays put relative to the
    Arabic-Indic digits around it.

    Also replace the thousands-separator "," BETWEEN DIGITS with the Arabic
    thousands separator "٬" (U+066C), same reasoning. Standalone commas in
    sentences (like "حي النرجس، الرياض") are left alone because the Arabic
    comma is already being used there.

    WHY convert digits at all: Arabic-Indic digits are classified as strong-RTL
    by the Unicode bidi algorithm, so they inherit the surrounding Arabic
    direction naturally. Western digits are strong-LTR and create direction
    conflicts inside Arabic paragraphs — conflicts that manifest as periods
    landing in the wrong place, % drifting, etc.

    We keep Latin letters (A, B, C) untouched because those are letters, not
    digits — building codes like "(A/B/C/D)" stay in Latin.
    """
    if not isinstance(text, str) or not text:
        return text
    # First: convert digits
    text = text.translate(_WESTERN_TO_ARABIC_INDIC)
    # Then: convert "." and "," that are BETWEEN TWO DIGITS (i.e. they're part
    # of a number) to their Arabic equivalents. The regex uses lookbehind and
    # lookahead on the Arabic-Indic digit range so we don't touch standalone
    # punctuation in sentences.
    text = _re.sub(r"(?<=[٠-٩])\.(?=[٠-٩])", "٫", text)
    text = _re.sub(r"(?<=[٠-٩]),(?=[٠-٩])", "٬", text)
    return text

def _clean_text(text):
    """Normalize spacing, then convert Western digits to Arabic-Indic.
    Apply to every user-facing string."""
    return _convert_digits(_normalize_spacing(text))

def _strip_hyperlink_style(run, color_rgb):
    """Defeat PowerPoint's auto-styling of hyperlink runs.

    The problem: when python-pptx adds a hyperlink via run.hyperlink.address,
    it emits an <a:hlinkClick> element that PowerPoint treats as a cue to
    apply the theme's `hlink` color (blue) AND an underline — overriding the
    run's own color and underline settings.

    Setting `u="none"` on rPr suppresses the underline SOMETIMES but not
    reliably in PowerPoint desktop. Setting run.font.color.rgb gets clobbered
    because the theme override runs after.

    The reliable fix is to inject a <a:solidFill> directly INSIDE the
    <a:hlinkClick> element. This is a standard OOXML trick: the hlinkClick
    element can carry its own fill, which PowerPoint respects over the theme.
    Same for an explicit <a:uFill> (underline fill) and the u="none" attribute.

    We do all three:
      - u="none" on rPr (kills underline attribute)
      - <a:uFill><a:noFill/></a:uFill> on rPr (kills underline even if theme re-adds)
      - Explicit solidFill INSIDE the hlinkClick (overrides theme hlink color)
    """
    rPr = run._r.get_or_add_rPr()
    hlink = rPr.find(qn("a:hlinkClick"))
    if hlink is None:
        return

    # 1. Explicit u="none" on the run
    rPr.set("u", "none")

    # 2. Remove any existing uFill and add <a:uFill><a:noFill/></a:uFill>
    for existing in rPr.findall(qn("a:uFill")):
        rPr.remove(existing)
    uFill = etree.SubElement(rPr, qn("a:uFill"))
    etree.SubElement(uFill, qn("a:noFill"))

    # 3. Inject a solidFill INSIDE the hlinkClick so PowerPoint uses our color
    #    instead of the theme's hlink color. The XML here is fragile — the
    #    solidFill must be a child of hlinkClick, not a sibling.
    for existing in hlink.findall(qn("a:solidFill")):
        hlink.remove(existing)
    solidFill = etree.SubElement(hlink, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    # color_rgb is an RGBColor (int subclass); convert to hex string
    srgb.set("val", str(color_rgb))

def _set_autofit_mode(tf, mode):
    """Set the text-frame's auto-fit behavior.

    Modes:
      - "shrink": `<a:normAutofit/>` — text shrinks to fit a fixed shape when
        it would otherwise overflow. Shape size is unchanged. Does nothing if
        text already fits (no effect on under-flow). Use for fixed-grid layouts
        where neighbor shapes can't be pushed out of place.
      - "grow":   `<a:spAutoFit/>` — shape resizes to match text exactly. Shape
        grows when text is long AND shrinks when text is short. Use for
        flow/stacked layouts. WARNING: shifts the bottom edge of the shape,
        which can collide with downstream shapes if the layout assumes fixed
        positions.
      - "none":   `<a:noAutofit/>` — neither; text can overflow the shape
        visually. Only use when you specifically want that.

    python-pptx's `auto_size` enum only covers "grow" (SHAPE_TO_FIT_TEXT). The
    "shrink" mode requires raw XML because the enum has no value for it.
    """
    assert mode in ("shrink", "grow", "none")
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = etree.SubElement(tf._txBody, qn("a:bodyPr"))
    # Remove any existing autofit mode; only one can be active.
    for tag in ("a:noAutofit", "a:spAutoFit", "a:normAutofit"):
        for existing in bodyPr.findall(qn(tag)):
            bodyPr.remove(existing)
    tag_map = {"shrink": "a:normAutofit", "grow": "a:spAutoFit", "none": "a:noAutofit"}
    etree.SubElement(bodyPr, qn(tag_map[mode]))

# Backward-compat alias — older code uses the shrink-only helper name.
def _enable_text_autofit(tf):
    _set_autofit_mode(tf, "shrink")

def _add_text(slide, x, y, w, h, text, *,
              font_size=12, bold=False, color=CHARCOAL,
              align="right", rtl=True, font_name=FONT,
              anchor="middle", hyperlink=None, shape_hyperlink=None,
              normalize=True, line_spacing=None,
              auto_fit=True, grow_to_fit=False):
    """Drop-in text box. Handles Amiri, color, RTL, alignment, optional hyperlink.

    `text` can be a string or a list of strings (each becomes one paragraph).
    Set `normalize=False` to bypass spacing/number cleanup.

    Line spacing rule: if the text block has 8 or more words total (summed
    across all paragraphs), every paragraph in the block gets 1.5 line spacing.
    Shorter blocks (titles, KPI labels, short bullets) stay at default single
    spacing, which keeps them visually tight.

    Pass `line_spacing` explicitly (e.g. 1.2 or 2.0) to override the auto rule.

    Fit behavior (two mutually-exclusive options):

      - `auto_fit=True` (DEFAULT) — text shrinks to fit the shape when content
        would overflow. Shape size stays fixed. Does nothing if text already
        fits (no effect on under-flow). This is the safer default for a
        fixed-grid deck like Wassel: it protects the layout if future content
        is longer than expected, at the cost of potentially smaller font.

      - `grow_to_fit=True` — shape resizes to match the text exactly. Grows
        when text is long, shrinks when text is short. Use for boxes that
        aren't in a fixed grid (e.g. the footer tagline, or any box where you
        genuinely want the shape to wrap the content). WARNING: this shifts
        the shape's bottom edge, so neighboring shapes below must be able to
        reflow or be tolerant of the movement. If both options are passed,
        `grow_to_fit` wins.

      - Pass `auto_fit=False, grow_to_fit=False` to opt out of both (text can
        overflow visually). Rarely needed.

    Hyperlink options:
      - `hyperlink`: run-level hyperlink. PowerPoint applies theme hlink color
        and underline, which overrides run.font.color.rgb. Avoid for styled links.
      - `shape_hyperlink`: shape-level hyperlink (click target = whole text box).
        PowerPoint does NOT apply theme hyperlink styling to shape-level links,
        so the text keeps its run color/underline settings intact. Preferred
        when you want a clickable but visually-styled text."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    # Fit behavior is ALWAYS explicitly set (never left to python-pptx defaults,
    # which use spAutoFit — that caused slide 4's insight card to grow into
    # its neighbour). Priority: grow_to_fit > auto_fit > none.
    if grow_to_fit:
        _set_autofit_mode(tf, "grow")
    elif auto_fit:
        _set_autofit_mode(tf, "shrink")
    else:
        _set_autofit_mode(tf, "none")
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map[anchor]

    lines = text if isinstance(text, list) else [text]

    # Determine line spacing: explicit override > auto rule.
    # Word count is across the whole block (all paragraphs combined).
    # We strip bidi marks before counting so LRM/ALM/RLM don't inflate the count.
    if line_spacing is None:
        def _strip_marks(s):
            return "".join(ch for ch in s if ch not in (LRM, RLM, ALM, LRE, PDF))
        total_words = sum(len(_strip_marks(ln).split()) for ln in lines)
        line_spacing = 1.5 if total_words >= 8 else None

    for i, ln in enumerate(lines):
        if normalize:
            ln = _clean_text(ln)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        p.alignment = align_map[align]
        if rtl:
            _set_rtl(p)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        _set_font_all_scripts(run, font_name)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        if hyperlink:
            run.hyperlink.address = hyperlink
            _strip_hyperlink_style(run, color)

    if shape_hyperlink:
        _add_shape_hyperlink(tb, shape_hyperlink)

    return tb


def _add_shape_hyperlink(shape, url):
    """Attach a hyperlink to the SHAPE (not a text run).

    PowerPoint renders shape-level hyperlinks as clickable but does NOT apply
    the theme's hlink color/underline to the text inside — those only trigger
    when the <a:hlinkClick> is inside a run's <a:rPr>. By placing it on the
    shape's <p:cNvPr> instead, we keep the click behavior and the text stays
    whatever color we set.

    Requires adding a relationship to the slide's rels file; python-pptx
    exposes this via the slide's part.
    """
    # Add a relationship from the slide part to the URL
    slide_part = shape.part
    rId = slide_part.relate_to(
        url,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )

    # Find the shape's <p:nvSpPr>/<p:cNvPr> and add <a:hlinkClick> as a child
    sp = shape._element  # <p:sp>
    nvSpPr = sp.find(qn("p:nvSpPr"))
    if nvSpPr is None:
        return
    cNvPr = nvSpPr.find(qn("p:cNvPr"))
    if cNvPr is None:
        return
    # Remove any existing hlinkClick so we don't double up
    for existing in cNvPr.findall(qn("a:hlinkClick")):
        cNvPr.remove(existing)
    hlink = etree.SubElement(cNvPr, qn("a:hlinkClick"))
    hlink.set(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id",
        rId,
    )

def _add_logo(slide, x, y, w, h):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(x), Inches(y),
                                  width=Inches(w), height=Inches(h))

def _get(content, *path, default="—"):
    """Safe nested lookup that returns a placeholder if anything in the path is missing."""
    cur = content
    for key in path:
        if cur is None or not isinstance(cur, dict) or key not in cur:
            return default
        cur = cur[key]
    return cur if cur not in ("", None) else default

# ─── Chrome helpers (named to match the Wassel build spec) ────────────────────

def addStrip(slide, *, variant="content"):
    """Left-edge ornamental strip.

    content variant: narrow sand strip + 9 small copper dots (minimal)
    divider variant: full copper strip with a Najdi/Diriyah-style ornament:
        a column of inward-pointing triangular notches (the classic stepped-triangle
        motif you see on Diriyah mud-brick walls) with a small diamond between
        each pair — not the dot ladder, which reads as generic.
    """
    if variant == "content":
        _add_rect(slide, 0.00, 0.00, 0.18, 5.625, SAND)
        for y in [0.50, 1.06, 1.62, 2.18, 2.74, 3.30, 3.86, 4.42, 4.98]:
            _add_rect(slide, 0.04, y, 0.10, 0.10, COPPER)
    elif variant == "divider":
        # Full-height copper strip
        _add_rect(slide, 0.00, 0.00, 0.55, 5.625, COPPER)

        # Najdi motif: triangular notches cut from the inner edge (pointing right
        # into the strip) + a small brown diamond between each pair of notches.
        # Draw notches as filled BROWN triangles so they "eat into" the copper.
        # Strip occupies x=0.00..0.55 (width 0.55). Notches at inner edge (x≈0.35..0.55).
        notch_count = 9
        strip_top, strip_bottom = 0.15, 5.47
        total_h = strip_bottom - strip_top
        gap = total_h / notch_count
        notch_h = 0.28  # triangle vertical span
        notch_depth = 0.20  # how far it cuts in (into the copper strip)

        for i in range(notch_count):
            cy = strip_top + gap * (i + 0.5)  # triangle center y
            # Filled brown triangle pointing from right edge (x=0.55) inward.
            # Points: right-top, tip (pointing left), right-bottom
            _add_right_pointing_notch(slide,
                                       right_x=0.55,
                                       tip_x=0.55 - notch_depth,
                                       top_y=cy - notch_h / 2,
                                       bot_y=cy + notch_h / 2)

            # Small brown diamond between this notch and the next (skip last)
            if i < notch_count - 1:
                diamond_cy = cy + gap / 2
                _add_small_diamond(slide,
                                    cx=0.14,
                                    cy=diamond_cy,
                                    size=0.10,
                                    fill=BROWN)


def _add_right_pointing_notch(slide, right_x, tip_x, top_y, bot_y):
    """A filled BROWN triangle pointing LEFT (tip at tip_x, base at right_x).

    Uses PowerPoint's built-in RIGHT_TRIANGLE preset, rotated -90° so the
    hypotenuse-free vertex (the tip) points left. This renders reliably in
    PowerPoint, WPS, LibreOffice, and Keynote — unlike custGeom which
    PowerPoint's stricter parser can reject silently.

    Implementation note: we'd ideally use ISOCELES_TRIANGLE rotated 270°, but
    python-pptx's ISOCELES_TRIANGLE rotates around its center, and the default
    orientation points up. 270° rotation gives us a left-pointing triangle
    with tip on the left side of the bounding box and a flat right edge —
    exactly what we want for a notch biting into the copper strip from its
    right edge.
    """
    bbox_left = min(tip_x, right_x)
    bbox_right = max(tip_x, right_x)
    bbox_w = bbox_right - bbox_left
    bbox_h = bot_y - top_y

    shp = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  Inches(bbox_left), Inches(top_y),
                                  Inches(bbox_w), Inches(bbox_h))
    # Default ISOCELES_TRIANGLE points up (tip at top-center). Rotate -90°
    # (= 270°) to make the tip point LEFT.
    shp.rotation = 270
    shp.fill.solid()
    shp.fill.fore_color.rgb = BROWN
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_small_diamond(slide, cx, cy, size, fill):
    """A small filled diamond (square rotated 45°) centered at (cx, cy)."""
    half = size / 2
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                  Inches(cx - half), Inches(cy - half),
                                  Inches(size), Inches(size))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def addHeader(slide, title, subtitle):
    """Brown header band with gold underline, copper right cap, logo, title and subtitle."""
    _add_rect(slide, 0.00, 0.00,  10.00, 1.10, BROWN)
    _add_rect(slide, 0.00, 1.10,  10.00, 0.04, GOLD)
    _add_rect(slide, 9.85, 0.00,   0.15, 1.10, COPPER)
    _add_logo(slide, 0.28, 0.16, 0.57, 0.57)
    _add_text(slide, 1.10, 0.10, 8.55, 0.56, title,
              font_size=26, bold=True, color=WHITE, align="right", rtl=True, anchor="middle")
    _add_text(slide, 1.10, 0.65, 8.55, 0.38, subtitle,
              font_size=13, color=GOLD, align="right", rtl=True, anchor="middle")

def addFooter(slide, right_text):
    """Sand footer band with copper top border. Right: project context. Left: wassel.re link."""
    _add_rect(slide, 0.00, 5.37, 10.00, 0.26, SAND)
    _add_rect(slide, 0.00, 5.37, 10.00, 0.02, COPPER)
    _add_text(slide, 0.25, 5.38, 9.50, 0.23, right_text,
              font_size=9, color=BROWN, align="right", rtl=True, anchor="middle")
    _add_text(slide, 0.25, 5.38, 9.50, 0.23, "wassel.re",
              font_size=9, color=COPPER, align="left", rtl=False, anchor="middle",
              shape_hyperlink="https://wassel.re")

def addDivider(slide, title, subtitle):
    """Full divider slide body: brown bg, left copper strip, centered title, gold line, subtitle."""
    _add_rect(slide, 0, 0, 10, 5.625, BROWN)
    addStrip(slide, variant="divider")
    _add_rect(slide, 1.00, 2.65, 8.50, 0.04, GOLD)
    _add_text(slide, 0.70, 1.55, 9.00, 1.05, title,
              font_size=52, bold=True, color=CREAM, align="center", rtl=True, anchor="middle")
    _add_text(slide, 0.70, 2.72, 9.00, 0.50, subtitle,
              font_size=20, color=GOLD, align="center", rtl=True, anchor="middle")

def _blank_slide(prs, bg_color=CREAM):
    """New slide with background."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # background rectangle covering the whole slide
    _add_rect(slide, 0, 0, 10, 5.625, bg_color)
    return slide

# ─── Per-slide builders ───────────────────────────────────────────────────────

def slide1_cover(prs, content):
    s = _blank_slide(prs, bg_color=BROWN)
    addStrip(s, variant="divider")
    _add_rect(s, 1.50, 2.75, 7.00, 0.03, GOLD)
    _add_logo(s, 3.94, 0.08, 2.56, 2.56)
    _add_text(s, 0.70, 2.85, 9.10, 0.85, "وصل العقارية",
              font_size=46, bold=True, color=CREAM, align="center", rtl=True)
    _add_text(s, 0.70, 3.65, 9.10, 0.48, "خطة التسويق والإدارة البيعية",
              font_size=21, color=SAND, align="center", rtl=True)
    # Order: project — city، district. City comes before district per brand spec.
    # This also helps bidi: the digit at the end of the project name (e.g., "17")
    # now sits immediately before " — " and the next Arabic word, with no neutral
    # characters able to reorder across it.
    tag = f"{_get(content, 'project_name')} — {_get(content, 'city')}، {_get(content, 'district')}"
    _add_rect(s, 2.80, 4.28, 4.40, 0.55, COPPER)
    _add_text(s, 2.80, 4.29, 4.40, 0.52, tag,
              font_size=13.5, bold=True, color=WHITE, align="center", rtl=True)
    _add_text(s, 0.70, 4.98, 9.10, 0.35, _get(content, 'year'),
              font_size=13, color=GOLD, align="center", rtl=False)

def slide2_about(prs, content):
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    addHeader(s, "عن وصل العقارية", "منظومة بيع عقاري متكاملة")
    addFooter(s, _get(content, 'footer_right'))

    # Body paragraph
    _add_text(s, 0.85, 1.17, 7.49, 0.62,
              ["شركة سعودية متخصصة في التسويق العقاري وإدارة المبيعات، تعمل على مستوى المملكة العربية السعودية.",
               "نؤمن بأن البيع الناجح يبدأ بنظام لا بجهد فردي."],
              font_size=15, color=BROWN, align="center", rtl=True)

    # 3 brown stat cards (right-to-left in Arabic: card 1 is physically leftmost)
    stats = [
        (1.38, "1 مليون+",    "ميزانيات مصروفة على الحملات"),
        (3.76, "70+",          "مشروع مدار"),
        (6.14, "2+ مليار",    "مشاريع مباعة"),
    ]
    for x, num, lbl in stats:
        _add_rect(s, x, 1.96, 2.20, 1.22, BROWN)
        _add_rect(s, x, 1.96, 2.20, 0.06, GOLD)
        _add_text(s, x, 2.06, 2.20, 0.62, num,
                  font_size=34, bold=True, color=GOLD, align="center", rtl=True)
        _add_text(s, x, 2.68, 2.20, 0.42, lbl,
                  font_size=12, color=SAND, align="center", rtl=True)

    # 3 white value-prop cards
    props = [
        (1.38, "عملية مبيعات موحدة", "ضمان عملية مبيعات عالية الكفاءة تعتمد على النظام لا الفرد"),
        (3.76, "نظام تقني متكامل",   "أتمتة كاملة من الاستقبال حتى الإفراغ"),
        (6.14, "تسويق ذكي مبني على البيانات", "محتوى وإعلانات مركّزة تصل للعميل الصحيح"),
    ]
    for x, title, body in props:
        _add_rect(s, x, 3.37, 2.20, 1.77, WHITE)
        _add_rect(s, x, 3.37, 2.20, 0.06, COPPER)
        _add_text(s, x + 0.10, 3.47, 2.00, 0.38, title,
                  font_size=13, bold=True, color=BROWN, align="center", rtl=True)
        _add_text(s, x + 0.10, 3.89, 2.00, 1.10, body,
                  font_size=11.5, color=CHARCOAL, align="center", rtl=True)

def slide3_divider(prs, content):
    s = _blank_slide(prs, bg_color=BROWN)
    addDivider(s, "تحليل مربع المشروع", "دراسة السوق والمشروع وفرصته")

def slide4_market(prs, content):
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    project = _get(content, 'project_name')
    district = _get(content, 'district')
    city = _get(content, 'city')
    # EXACT pattern required by the spec: "مربع مشروع <X> <Y>"
    subtitle = f"تحليل مربع مشروع {project} {district} — {city}"
    addHeader(s, "تحليل السوق", subtitle)
    addFooter(s, _get(content, 'footer_right'))

    # KPI tiles are CONTENT-DRIVEN, not fixed-slot. The research+merge layer
    # picks the three most presentable/logical metrics from what actually got
    # filled in sources.csv for THIS project. No hardcoded indicator names.
    # Schema: content["market"]["kpis"] = [{"value": "10,202", "label": "متوسط سعر متر الدور في الحي"}, ...]
    # Exactly 3 items. Each label should be self-explanatory about scope
    # (e.g. "...في الحي" for district-level, "...للمشروع" for project-level)
    # so the viewer can't misread a حي figure as a project figure.
    xs = [0.30, 3.40, 6.50]
    kpis = _get(content, 'market', 'kpis', default=[])
    if not isinstance(kpis, list) or len(kpis) == 0:
        kpis = [{"value": "—", "label": ""}] * 3
    kpis = (kpis + [{"value": "—", "label": ""}] * 3)[:3]
    for x, kpi in zip(xs, kpis):
        num = kpi.get("value", "—") if isinstance(kpi, dict) else "—"
        lbl = kpi.get("label", "") if isinstance(kpi, dict) else ""
        _add_rect(s, x, 1.25, 2.90, 1.65, WHITE)
        _add_rect(s, x, 1.25, 2.90, 0.06, COPPER)
        _add_text(s, x, 1.38, 2.90, 0.75, num,
                  font_size=42, bold=True, color=COPPER, align="center", rtl=True)
        _add_text(s, x, 2.12, 2.90, 0.34, lbl,
                  font_size=12, color=CHARCOAL, align="center", rtl=True)

    # Insight card.
    # HARD SPEC: each insight line should be ≤ 12 words (ideally ≤ 8) and there
    # should be exactly 2 lines. The card is sized for punchy one-liners, not
    # full sentences. `auto_fit=True` shrinks the font if content is slightly
    # over budget, but if insights are 2x too long, the content layer (research
    # / presentation skill) must trim them — not the renderer.
    _add_rect(s, 0.30, 3.08, 9.10, 1.10, WHITE)
    _add_rect(s, 9.41, 3.08, 0.08, 1.10, COPPER)
    _add_text(s, 0.28, 3.35, 0.25, 0.45, "✦",
              font_size=18, color=COPPER, align="center", rtl=False)
    insight_lines = _get(content, 'market', 'insight_lines', default=["—", ""])
    if isinstance(insight_lines, str):
        insight_lines = [insight_lines]
    _add_text(s, 0.55, 3.12, 8.85, 0.98, insight_lines,
              font_size=14, color=BROWN, align="right", rtl=True,
              auto_fit=True)

    # Price-range row
    _add_rect(s, 0.30, 4.32, 9.19, 0.78, WHITE)
    _add_rect(s, 5.90, 4.36, 0.02, 0.68, SAND)
    rmin = _get(content, 'market', 'price_range_min')
    rmax = _get(content, 'market', 'price_range_max')
    avg  = _get(content, 'market', 'avg_price')
    _add_text(s, 0.35, 4.38, 5.40, 0.28, "نطاق الأسعار في المربع:",
              font_size=11, bold=True, color=BROWN, align="right", rtl=True)
    _add_text(s, 0.35, 4.62, 5.40, 0.28, f"{rmin} ريال الى {rmax} ريال",
              font_size=13, bold=True, color=COPPER, align="right", rtl=True)
    _add_text(s, 6.30, 4.38, 3.10, 0.28, "متوسط سعر الوحدة:",
              font_size=11, bold=True, color=BROWN, align="right", rtl=True)
    _add_text(s, 6.40, 4.61, 3.10, 0.28, f"~{avg} ريال",
              font_size=13, bold=True, color=COPPER, align="right", rtl=True)

def slide5_competitors(prs, content):
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    addHeader(s, "مقارنة المشاريع المنافسة",
              f"المشاريع السكنية النشطة في مربع {_get(content,'project_name')} {_get(content,'district')}")
    addFooter(s, _get(content, 'footer_right'))

    headers = _get(content, 'competitors', 'headers',
                   default=["المشروع", "المطوّر", "عدد الوحدات", "نسبة البيع", "متوسط السعر", "ملاحظات"])
    rows = _get(content, 'competitors', 'rows', default=[])
    if not rows:
        rows = [["—"] * len(headers)]

    # RTL HACK: PowerPoint tables don't have a true "right-to-left" flag that
    # reverses column order. The fix is to reverse the *data* so the first
    # logical column (in Arabic: rightmost visually) ends up as column 0 —
    # which PowerPoint paints on the physical left — but since Arabic is read
    # right-to-left, that physical-left position IS the last column in reading
    # order, and so on. Net effect: column 0 of the data ends up where an
    # Arabic reader expects to start reading.
    headers_rtl = list(reversed(headers))
    rows_rtl = [list(reversed(r)) for r in rows]

    table_shape = s.shapes.add_table(
        rows=len(rows_rtl) + 1, cols=len(headers_rtl),
        left=Inches(0.25), top=Inches(1.20),
        width=Inches(9.55), height=Inches(3.15)
    )
    tbl = table_shape.table
    # Header row
    for c, h in enumerate(headers_rtl):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BROWN
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_rtl(p)
        r = p.add_run()
        r.text = _clean_text(h)
        _set_font_all_scripts(r, FONT)
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE
    # Body rows
    for ri, row in enumerate(rows_rtl):
        for c, val in enumerate(row):
            cell = tbl.cell(ri + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SAND if ri % 2 == 0 else WHITE
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _set_rtl(p)
            r = p.add_run()
            r.text = _clean_text(str(val))
            _set_font_all_scripts(r, FONT)
            r.font.size = Pt(11)
            r.font.color.rgb = CHARCOAL

    footnote = _get(content, 'competitors', 'footnote', default="")
    if footnote:
        _add_text(s, 0.25, 5.12, 9.50, 0.22, footnote,
                  font_size=9.5, color=CHARCOAL, align="right", rtl=True)

def slide6_opportunity(prs, content):
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    addHeader(s, "المشروع وفرصته في السوق",
              f"{_get(content,'project_name')} — {_get(content,'district')}")
    addFooter(s, _get(content, 'footer_right'))

    # RIGHT: specs card
    _add_rect(s, 4.85, 1.22, 4.90, 3.85, WHITE)
    _add_rect(s, 4.85, 1.22, 4.90, 0.06, COPPER)
    _add_text(s, 4.90, 1.32, 4.80, 0.40, "مواصفات المشروع",
              font_size=15, bold=True, color=BROWN, align="right", rtl=True)
    rows = [
        ("الموقع",          _get(content, 'project', 'location')),
        ("المساحات",         _get(content, 'project', 'area_range')),
        ("الأسعار",          _get(content, 'project', 'price_range')),
        ("الوحدات الكلية",   _get(content, 'project', 'unit_total')),
        ("الوحدات المتاحة",  _get(content, 'project', 'units_available')),
        ("المرافق",          _get(content, 'project', 'amenities')),
        ("الضمانات",         _get(content, 'project', 'warranties')),
    ]
    y = 1.80
    for i, (label, value) in enumerate(rows):
        bg = CREAM if i % 2 == 0 else WHITE
        _add_rect(s, 4.90, y, 4.80, 0.40, bg)
        _add_text(s, 7.60, y + 0.05, 2.05, 0.30, label,
                  font_size=10, bold=True, color=COPPER, align="right", rtl=True)
        _add_text(s, 4.92, y + 0.05, 2.60, 0.30, value,
                  font_size=10, color=CHARCOAL, align="right", rtl=True)
        y += 0.45

    # LEFT TOP: brown opportunity card (bullets)
    _add_rect(s, 0.25, 1.22, 4.40, 1.88, BROWN)
    _add_rect(s, 0.25, 1.22, 4.40, 0.06, GOLD)
    _add_text(s, 0.30, 1.30, 4.25, 0.40, "الفرصة الاستراتيجية",
              font_size=14, bold=True, color=GOLD, align="right", rtl=True)
    bullets = _get(content, 'project', 'opportunity_bullets', default=["—"] * 4)
    while len(bullets) < 4:
        bullets.append("—")
    bullet_ys = [1.77, 2.07, 2.37, 2.67]
    for by, text in zip(bullet_ys, bullets[:4]):
        _add_text(s, 0.30, by, 4.30, 0.27, f"✦  {text}",
                  font_size=12, color=SAND, align="right", rtl=True)

    # LEFT BOTTOM: revenue card
    _add_rect(s, 0.25, 3.20, 4.40, 1.87, WHITE)
    _add_rect(s, 0.25, 3.20, 4.40, 0.06, COPPER)
    _add_text(s, 0.30, 3.30, 4.25, 0.38, "الإيراد المتوقع للمطوّر من الوحدات المتبقية",
              font_size=13, bold=True, color=BROWN, align="right", rtl=True)
    _add_text(s, 0.30, 3.68, 4.25, 0.72, _get(content, 'project', 'expected_revenue'),
              font_size=38, bold=True, color=COPPER, align="right", rtl=True)
    _add_text(s, 0.30, 4.38, 4.25, 0.26, "ريال سعودي  — تقديري",
              font_size=11, color=CHARCOAL, align="right", rtl=True)
    remaining = _get(content, 'project', 'remaining_units')
    avg_m     = _get(content, 'project', 'avg_price_millions')
    _add_text(s, 0.30, 4.66, 4.25, 0.32,
              f"{remaining} وحدة متبقية  ×  متوسط سعر {avg_m} مليون ريال",
              font_size=10.5, color=COPPER, align="right", rtl=True)

def slide7_divider(prs, content):
    s = _blank_slide(prs, bg_color=BROWN)
    # EXACT subtitle required by spec — do not edit
    addDivider(s, "الخطة التسويقية", "الهدف: صناعة الطلب، وجلب المهتمين")

def slide8_event(prs, content):
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    addHeader(s, "الخطة التسويقية – اولاً: حفل الإفتتاح",
              "إطلاق قوي يصنع الزخم من اليوم الأول")
    addFooter(s, _get(content, 'marketing_footer_right',
                      default=_get(content, 'footer_right')))

    # Invitees bar
    _add_rect(s, 0.41, 1.96, 9.16, 0.32, BROWN)
    _add_text(s, 0.38, 1.98, 9.30, 0.28, "المدعوون",
              font_size=11, color=GOLD, align="center", rtl=True)

    # 5 audience pills (right-to-left: rightmost is first in Arabic reading).
    # Icons are typographic symbols (not emoji) so they take our brand color.
    audience = [
        (0.41, COPPER, "◆", "دعوات عامة"),
        (2.26, BROWN,  "★", "عملاء وصل"),
        (4.11, COPPER, "♦", "مؤثرون عقاريون"),
        (5.96, BROWN,  "■", "مطوّرون ومستثمرون"),
        (7.81, COPPER, "●", "أهل الحي"),
    ]
    for x, fill, icon, label in audience:
        _add_rect(s, x, 2.33, 1.76, 0.50, fill)
        _add_text(s, x + 0.05, 2.36, 1.20, 0.46, label,
                  font_size=10.5, bold=True, color=WHITE, align="center", rtl=True)

    # Elements bar (3 equal cards)
    _add_rect(s, 0.41, 2.88, 9.16, 0.32, BROWN)
    _add_text(s, 0.38, 2.90, 9.30, 0.28, "عناصر الحفل",
              font_size=11, color=GOLD, align="center", rtl=True)

    # Event element cards. Icons are typographic symbols in COPPER (brand color)
    # instead of emoji (which render in their own colors and clash with palette).
    elements = _get(content, "event_elements", default=[
        ("♪", "عرض فني"),           # music note — performance
        ("◆", "لوحات إرشادية"),      # diamond — signage/wayfinding
        ("✦", "إنارة احترافية"),     # star — lighting
    ])
    n = len(elements)
    total_w = 9.16
    gap = 0.10
    card_w = (total_w - gap * (n - 1)) / n
    for i, (icon, label) in enumerate(elements):
        x = 0.41 + i * (card_w + gap)
        _add_rect(s, x, 3.23, card_w, 1.80, WHITE)
        _add_rect(s, x, 5.00, card_w, 0.04, GOLD)
        _add_text(s, x, 3.40, card_w, 0.80, icon,
                  font_size=60, bold=True, color=COPPER, align="center", rtl=False)
        _add_text(s, x, 4.30, card_w, 0.50, label,
                  font_size=14, bold=True, color=BROWN, align="center", rtl=True)

def slide9_content_platforms(prs, content):
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    addHeader(s, "الخطة التسويقية — ثانياً: المحتوى والتسويق الرقمي",
              "تسويق رقمي مركّز يصل للعميل الصحيح")
    addFooter(s, _get(content, 'marketing_footer_right',
                      default=_get(content, 'footer_right')))

    # "Content" bar
    _add_rect(s, 0.33, 1.22, 9.15, 0.30, BROWN)
    _add_text(s, 0.30, 1.24, 9.35, 0.26, "المحتوى",
              font_size=11, color=GOLD, align="center", rtl=True)

    # 6 content tiles in 3×2 grid (right-to-left reading).
    # Each tile: icon on the physically-LEFT side (in copper, brand color),
    # title + subtitle on the physically-RIGHT side (after copper accent bar).
    # Icons are typographic symbols — chosen to evoke each content type while
    # taking the font color (unlike emoji which ignore color):
    #   ▶ (play triangle)   — video / cinematic
    #   ◐ (half-filled circle) — virtual tour / 360°
    #   ✈ (airplane)        — drone / aerial
    #   ≋ (wave)            — promotional clips / social virality
    #   ♯ (sharp sign)      — social posts / hashtag
    #   ❐ (empty document)  — brochure
    tiles_top = [
        (6.53, "✈", "درون جوي",        "إبراز الموقع والمحيط"),
        (3.43, "◐", "جولات افتراضية",  "تقليل التردد قبل الزيارة"),
        (0.33, "▶", "فيديو سينمائي",   "إبراز الفخامة والتفاصيل"),
    ]
    tiles_bot = [
        (6.53, "≋", "مقاطع ترويجية",         "الانتشار الواسع والسريع"),
        (3.43, "♯", "منشورات منصات التواصل", "عروض ومزايا ولقطات"),
        (0.33, "❐", "بروشور رسمي",           "كتيب تعريفي عن المشروع ويتضمن قائمة الوحدات"),
    ]
    for y, row in [(1.63, tiles_top), (2.33, tiles_bot)]:
        for x, icon, title, sub in row:
            _add_rect(s, x, y, 2.95, 0.60, WHITE)
            # Right-edge copper accent bar (physically right)
            _add_rect(s, x + 2.86, y, 0.09, 0.60, COPPER)
            # Icon on the physically-LEFT side of the card, in copper
            _add_text(s, x + 0.05, y, 0.50, 0.60, icon,
                      font_size=22, bold=True, color=COPPER,
                      align="center", rtl=False, anchor="middle")
            # Title + subtitle, right-aligned, taking the remaining right portion
            _add_text(s, x + 0.55, y + 0.04, 2.26, 0.28, title,
                      font_size=12, bold=True, color=BROWN, align="right", rtl=True)
            _add_text(s, x + 0.55, y + 0.31, 2.26, 0.24, sub,
                      font_size=10, color=CHARCOAL, align="right", rtl=True)

    # Platforms bar
    _add_rect(s, 0.33, 3.14, 9.15, 0.28, BROWN)
    _add_text(s, 0.30, 3.16, 9.35, 0.23, "المنصات الترويجية",
              font_size=11, bold=True, color=GOLD, align="center", rtl=True)

    # 4 platform cards (right-to-left)
    platforms = [
        (0.33, RGBColor(0xFF, 0xFC, 0x00), "سناب شات",  "التحويل الأساسي",         "تغطيات للوحدات",                      BROWN),
        (2.65, RGBColor(0x00, 0x00, 0x00), "تيك توك",    "صناعة الطلب والانتشار",  "محتوى بصري سريع + جولات طبيعية + مؤثرون", WHITE),
        (5.05, RGBColor(0xC1, 0x35, 0x84), "إنستاغرام",  "بناء الثقة",               "ريلز احترافية + ستوري يومي + مراجعات",     WHITE),
        (7.45, RGBColor(0x0A, 0x66, 0xC2), "لينكدإن",     "استقطاب المستثمرين",     "عائد إيجاري + نمو القيمة + موقع استراتيجي", WHITE),
    ]
    for x, header_fill, name, role, body, name_color in platforms:
        w = 2.03 if x != 0.33 else 1.93
        _add_rect(s, x, 3.48, w, 1.77, WHITE)
        _add_rect(s, x, 3.48, w, 0.52, header_fill)
        _add_text(s, x, 3.52, w, 0.23, name,
                  font_size=13, bold=True, color=name_color, align="center", rtl=True)
        _add_text(s, x + 0.10, 4.13, w - 0.20, 0.24, role,
                  font_size=11, bold=True, color=COPPER, align="center", rtl=True)
        _add_text(s, x + 0.10, 4.42, w - 0.20, 0.84, body,
                  font_size=10, color=CHARCOAL, align="center", rtl=True)

def slide10_outcomes(prs, content):
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    addHeader(s, "الخطة التسويقية - ثالثاً: نتائج مدروسة",
              "قوة الخطة تكمن في تزامن المحاور الثلاثة")
    addFooter(s, _get(content, 'marketing_footer_right',
                      default=_get(content, 'footer_right')))

    # LEFT — target numbers card.
    # Values are DERIVED from the project's total units via FIXED conversion
    # constants (view→lead = 1%, lead→sale = 0.6%). Never override these;
    # the whole point of this slide is to show the funnel math working from
    # the project's sales goal backward.
    _add_rect(s, 0.25, 1.22, 4.50, 3.60, WHITE)
    _add_rect(s, 0.25, 1.22, 4.50, 0.06, COPPER)
    _add_text(s, 0.30, 1.32, 4.35, 0.38, "أرقام الترويج المستهدفة",
              font_size=14, bold=True, color=BROWN, align="right", rtl=True)

    # Fixed conversion constants for the marketing funnel
    VIEW_TO_LEAD = 0.01   # 1%
    LEAD_TO_SALE = 0.006  # 0.6%

    # Primary input: total units to sell (entire project, monthly × months)
    try:
        total_sold = int(str(_get(content, 'marketing_targets',
                                    'total_units_sold', default="0")).strip() or 0)
    except (ValueError, TypeError):
        total_sold = 0

    # Derive counts (round to nearest integer since we can't have fractional people)
    leads_target   = round(total_sold / LEAD_TO_SALE) if total_sold else 0
    impressions    = round(leads_target / VIEW_TO_LEAD) if leads_target else 0

    # Format with Western thousands separator (digit converter will handle the
    # rest; our _convert_digits replaces "," between digits with ٬ automatically)
    def _fmt_int(n):
        return f"{n:,}"

    metrics = [
        (_fmt_int(impressions),   "مشاهدة مستهدفة"),
        ("1%",                     "نسبة التحويل إلى مهتمين"),
        (_fmt_int(leads_target),  "مهتم مستهدف"),
        ("0.6%",                   "نسبة تحويل المهتمين إلى مبيعات"),
        (_fmt_int(total_sold),    "وحدة مباعة — إجمالي المشروع"),
    ]
    y = 1.80
    for i, (num, lbl) in enumerate(metrics):
        bg = RGBColor(0xFF, 0xF7, 0xF2) if i % 2 == 0 else WHITE
        _add_rect(s, 0.30, y, 4.35, 0.52, bg)
        _add_text(s, 0.32, y + 0.08, 1.50, 0.36, num,
                  font_size=22, bold=True, color=COPPER, align="center", rtl=False)
        _add_text(s, 1.88, y + 0.12, 2.72, 0.28, lbl,
                  font_size=12, color=CHARCOAL, align="right", rtl=True)
        y += 0.58

    # RIGHT — why it works (dark)
    _add_rect(s, 4.95, 1.22, 4.80, 3.60, BROWN)
    _add_rect(s, 4.95, 1.22, 4.80, 0.06, GOLD)
    _add_text(s, 5.00, 1.30, 4.65, 0.42, "لماذا تنجح هذه الخطة؟",
              font_size=14, bold=True, color=GOLD, align="right", rtl=True)

    items = [
        ("المحتوى يخلق الرغبة",       "فيديوهات تُحوّل المشروع إلى تجربة إستثنائية"),
        ("الترويج يجلب العميل",       "استهداف دقيق يصل للمشتري الجاهز"),
        ("الحفل يسرّع القرار",        "بيئة مثالية للحجز الفوري من اليوم الأول"),
        ("إدارة العملاء تُغلق الصفقة", "متابعة منهجية حتى اكتمال البيع"),
    ]
    y_positions = [1.91, 2.65, 3.39, 4.13]
    for y, (title, body) in zip(y_positions, items):
        _add_text(s, 5.05, y, 0.30, 0.28, "◈",
                  font_size=14, color=GOLD, align="center", rtl=False)
        _add_text(s, 5.35, y - 0.04, 4.35, 0.28, title,
                  font_size=12.5, bold=True, color=WHITE, align="right", rtl=True)
        _add_text(s, 5.35, y + 0.23, 4.35, 0.28, body,
                  font_size=11, color=SAND, align="right", rtl=True)

    # Tagline strip
    _add_rect(s, 0.25, 4.95, 9.50, 0.40, COPPER)
    _add_text(s, 0.30, 4.98, 9.35, 0.32, "نحن لا ننتظر العميل —  نحن نصنعه",
              font_size=15, bold=True, color=WHITE, align="center", rtl=True)

def slide11_divider(prs, content):
    s = _blank_slide(prs, bg_color=BROWN)
    # EXACT subtitle required
    addDivider(s, "الخطة البيعية", "تحويل الطلب والاهتمام إلى مبيعات")

def slide12_sales_snake(prs, content):
    """Sales journey — 10-stage horizontal serpent (user-designed reference).

    Layout: 2 rows × 5 cards.
      Row 1 (top): stages 1→5, physically R→L in Arabic reading.
      Row 2 (bottom): stages 6→10, physically L→R.
    Each card: clean rectangle + dark-brown number strip on physical-right edge.
    Connectors: thin gold bars between same-row cards; L-bridge from row1 to row2.
    Stage 10 (finale): full copper-filled card with white checkmark.
    """
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    addHeader(s, "عملية المبيعات —  مسار العمل",
              "رحلة العميل من أول تواصل حتى الإفراغ")
    addFooter(s, _get(content, 'sales_footer_right',
                      default=_get(content, 'footer_right')))

    card_w = 1.72
    card_h = 0.62
    gap    = 0.18
    num_strip_w = 0.34
    body_w = card_w - num_strip_w

    row1_y = 1.95
    row2_y = 3.45
    row_total_w = 5 * card_w + 4 * gap
    row_start_left = (10.0 - row_total_w) / 2   # center horizontally

    def row1_card_x(stage):
        """Stage 1-5 in row 1 (reads R→L). Stage 1 = physically rightmost."""
        phys_col = 5 - stage
        return row_start_left + phys_col * (card_w + gap)

    def row2_card_x(stage):
        """Stage 6-10 in row 2 (reads L→R). Stage 6 = physically leftmost."""
        phys_col = stage - 6
        return row_start_left + phys_col * (card_w + gap)

    titles = {
        1: "استقبال العميل",
        2: "التواصل وجمع البيانات",
        3: "حجز الموعد",
        4: "تأكيد الموعد",
        5: "الزيارة",
        6: "متابعة بعد الزيارة",
        7: "تقديم عرض السعر",
        8: "الحجز",
        9: "التمويل",
        10: "الإفراغ",
    }

    def draw_card(x, y, stage_num, title, body_fill=WHITE):
        """Draw a standard journey card with right-edge dark brown number strip."""
        strip_x = x + body_w
        _add_rect(s, strip_x, y, num_strip_w, card_h, BROWN)
        _add_text(s, strip_x, y, num_strip_w, card_h, f"{stage_num:02d}",
                  font_size=13, bold=True, color=GOLD,
                  align="center", rtl=False, anchor="middle")
        _add_rect(s, x, y, body_w, card_h, body_fill, line_rgb=SAND)
        _add_text(s, x + 0.05, y, body_w - 0.10, card_h, title,
                  font_size=11, bold=True, color=BROWN,
                  align="center", rtl=True, anchor="middle")

    # ─── Connectors ────────────────────────────────────────────────────────
    connector_w = 0.06
    connector_y1 = row1_y + card_h / 2 - connector_w / 2
    connector_y2 = row2_y + card_h / 2 - connector_w / 2

    # Row 1 horizontal connectors (stage s to s+1: s is physically RIGHT of s+1)
    for stage in range(1, 5):
        xl = row1_card_x(stage + 1) + card_w
        xr = row1_card_x(stage)
        _add_rect(s, xl, connector_y1, xr - xl, connector_w, GOLD)
    # Row 2 horizontal connectors (s is physically LEFT of s+1)
    for stage in range(6, 10):
        xl = row2_card_x(stage) + card_w
        xr = row2_card_x(stage + 1)
        _add_rect(s, xl, connector_y2, xr - xl, connector_w, GOLD)

    # Row1→Row2 L-bridge: from physical-left of row 1 (stage 5) down to row 2 (stage 6).
    # Stages 5 and 6 are both at the leftmost physical position, so bridge is vertical.
    bridge_x = row1_card_x(5)
    _add_rect(s, bridge_x, row1_y + card_h,
               connector_w, row2_y - (row1_y + card_h), GOLD)
    # Short stub extending right into stage 6's card
    _add_rect(s, bridge_x, row2_y - connector_w / 2,
               0.22, connector_w, GOLD)

    # ─── Draw cards ────────────────────────────────────────────────────────
    for stage in range(1, 6):
        draw_card(row1_card_x(stage), row1_y, stage, titles[stage])
    for stage in range(6, 10):
        draw_card(row2_card_x(stage), row2_y, stage, titles[stage])

    # Stage 10 finale — full copper card
    x10 = row2_card_x(10)
    _add_rect(s, x10 + body_w, row2_y, num_strip_w, card_h, BROWN)
    _add_text(s, x10 + body_w, row2_y, num_strip_w, card_h, "10",
              font_size=13, bold=True, color=GOLD,
              align="center", rtl=False, anchor="middle")
    _add_rect(s, x10, row2_y, body_w, card_h, COPPER)
    check_w = 0.40
    _add_text(s, x10 + 0.05, row2_y, check_w, card_h, "✓",
              font_size=22, bold=True, color=WHITE,
              align="center", rtl=False, anchor="middle")
    label_x = x10 + check_w + 0.05
    label_w = body_w - check_w - 0.10
    _add_text(s, label_x, row2_y, label_w, card_h, titles[10],
              font_size=14, bold=True, color=WHITE,
              align="center", rtl=True, anchor="middle")


def slide_shapes_add_oval(slide, x, y, w, h, fill_rgb, line_rgb=None):
    """Helper to add a filled oval (circle if w==h) with optional outline."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp




def slide13_sales_detailed(prs, content):
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    addHeader(s, "عملية المبيعات —  النسخة التفصيلية",
              "10 خطوات منهجية من أول تواصل حتى إغلاق الصفقة")
    addFooter(s, _get(content, 'sales_footer_right',
                      default=_get(content, 'footer_right')))

    steps = [
        ("01", "استقبال العميل",        "تسجيل تلقائي في النظام وبدء مسار المتابعة فوراً"),
        ("02", "التواصل وجمع البيانات", "اتصال خلال دقائق، جمع نوع الوحدة والميزانية والهدف، تصنيف الجدية"),
        ("03", "حجز الموعد",             "اقتراح وقت محدد، تثبيته في النظام، إرسال تأكيد فوري للعميل"),
        ("04", "تأكيد الموعد",           "تذكير قبل يوم واتصال في نفس اليوم لرفع نسبة الحضور"),
        ("05", "الزيارة الميدانية",       "استقبال العميل وعرض الوحدات المناسبة بتجربة احترافية"),
        ("06", "المتابعة بعد الزيارة",    "تواصل فوري، معالجة الاعتراضات، إعادة تقديم العرض"),
        ("07", "تقديم عرض السعر",         "عرض مخصص بحسب الوحدة والميزانية يُرسَل بشكل منظم"),
        ("08", "الحجز",                   "تثبيت الوحدة وإنشاء مستند الحجز وتحديث حالتها في النظام"),
        ("09", "التمويل",                 "متابعة إجراءات التمويل والتنسيق مع الجهات المعنية"),
        ("10", "الإفراغ",                 "استكمال الإجراءات النهائية ونقل الملكية وإصدار المستندات"),
    ]
    # Right column (physically right): steps 1-5  -> x=5.08, badge x=9.60
    # Left column (physically left): steps 6-10   -> x=0.25, badge x=4.77
    ys = [1.22, 2.04, 2.86, 3.68, 4.50]

    for i, (num, title, body) in enumerate(steps[:5]):
        y = ys[i]
        bg = WHITE if i % 2 == 0 else CREAM
        _add_rect(s, 5.08, y, 4.68, 0.74, bg)
        _add_rect(s, 9.60, y + 0.17, 0.35, 0.35, COPPER)
        _add_text(s, 9.60, y + 0.17, 0.35, 0.35, num,
                  font_size=11, bold=True, color=WHITE, align="center", rtl=False)
        _add_text(s, 5.18, y + 0.06, 4.35, 0.28, title,
                  font_size=12, bold=True, color=BROWN, align="right", rtl=True)
        _add_text(s, 5.18, y + 0.37, 4.35, 0.30, body,
                  font_size=9.5, color=CHARCOAL, align="right", rtl=True)

    for i, (num, title, body) in enumerate(steps[5:]):
        y = ys[i]
        bg = WHITE if i % 2 == 0 else CREAM
        _add_rect(s, 0.25, y, 4.68, 0.74, bg)
        _add_rect(s, 4.77, y + 0.17, 0.35, 0.35, COPPER)
        _add_text(s, 4.77, y + 0.17, 0.35, 0.35, num,
                  font_size=11, bold=True, color=WHITE, align="center", rtl=False)
        _add_text(s, 0.35, y + 0.06, 4.35, 0.28, title,
                  font_size=12, bold=True, color=BROWN, align="right", rtl=True)
        _add_text(s, 0.35, y + 0.37, 4.35, 0.30, body,
                  font_size=9.5, color=CHARCOAL, align="right", rtl=True)

def slide14_funnel(prs, content):
    """Sales plan by the numbers — a formula-driven breakdown table.

    This slide was previously a 5-level funnel; it's now a detailed metrics
    table because a full row-by-row breakdown communicates the sales plan
    more precisely than a visual funnel (and the funnel trapezoids were
    rendering inconsistently in PowerPoint).

    FIXED formula constants (do NOT override via content dict):
      appointment_booking_rate = 6%     leads × 6%  → appointments
      appointment_attendance   = 40%    appts × 40% → appointment visits
      natural_visits_multiplier = 2     natural = 2 × appointment visits
      interested_from_visits    = 20%   total_visits × 20% → offers
      booking_rate              = 60%   offers × 60% → bookings
      sale_from_booking         = 80%   bookings × 80% → sales

    Primary inputs:
      - `sales_plan.leads_per_month` — monthly interested-leads baseline
      - `marketing_targets.total_units_sold` — TOTAL campaign units (not monthly).
        Used to decide monthly-distribute vs collapse-to-campaign presentation.

    ── Two presentation modes ──────────────────────────────────────────────
    MODE A (distribute normally):
        When total_units >= 3 AND natural_monthly_sales >= 4, the table shows
        a realistic monthly row. Header reads "الخطة البيعية الشهرية".
        The last-row sales number is the monthly figure.

    MODE B (collapse to single campaign):
        When total_units >= 3 AND natural_monthly_sales < 4 — OR when
        total_units < 3 (small projects can never hit 4/month) — the table
        collapses to a single "whole campaign" walk: leads are scaled so the
        derived sales equal total_units, and the row previously labeled
        "المبيعات" is re-labeled "إجمالي المبيعات". Header reads
        "الخطة البيعية للحملة".

        The COLLAPSE is a presentation choice for the table only. It does NOT
        lie about pacing. The sidebar KPI "أشهر لبيع المتبقي من المشروع" is
        still derived from the real natural_monthly_sales rate (honest
        schedule), independent of which mode the table is in.

    Rounding: round to nearest integer (round()) at every step.
    Percentages display with 2 decimals (6.00%, 40.00%) except derived overall
    rates which use 1 decimal.
    """
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    addFooter(s, _get(content, 'sales_footer_right',
                      default=_get(content, 'footer_right')))

    # ─── Formula constants (brand — do not override) ───────────────────────
    APPT_BOOKING_RATE   = 0.06   # 6%
    APPT_ATTENDANCE     = 0.40   # 40%
    NATURAL_VISITS_MULT = 2      # natural = 2 × appointment visits
    INTERESTED_FROM_VIS = 0.20   # 20%
    BOOKING_RATE        = 0.60   # 60%
    SALE_FROM_BOOKING   = 0.80   # 80%
    # Product of all steps = sales per 1 lead (before rounding cascade)
    SALE_PER_LEAD = (APPT_BOOKING_RATE * APPT_ATTENDANCE
                     * (1 + NATURAL_VISITS_MULT) * INTERESTED_FROM_VIS
                     * BOOKING_RATE * SALE_FROM_BOOKING)

    def _walk(n_leads):
        """Run the funnel forward with rounding. Returns dict of stage counts."""
        appts  = round(n_leads * APPT_BOOKING_RATE)
        a_vis  = round(appts * APPT_ATTENDANCE)
        nat    = NATURAL_VISITS_MULT * a_vis
        t_vis  = a_vis + nat
        offrs  = round(t_vis * INTERESTED_FROM_VIS)
        books  = round(offrs * BOOKING_RATE)
        sls    = round(books * SALE_FROM_BOOKING)
        return {"leads": n_leads, "appts": appts, "a_vis": a_vis, "nat": nat,
                "t_vis": t_vis, "offers": offrs, "books": books, "sales": sls}

    # ─── Read inputs ────────────────────────────────────────────────────────
    try:
        monthly_leads = int(str(_get(content, 'sales_plan', 'leads_per_month',
                                     default="0")).strip() or 0)
    except (ValueError, TypeError):
        monthly_leads = 0
    try:
        total_units = int(str(_get(content, 'marketing_targets',
                                   'total_units_sold', default="0")).strip() or 0)
    except (ValueError, TypeError):
        total_units = 0

    # ─── Natural (monthly) walk — basis for sidebar KPIs regardless of mode ─
    nat_walk = _walk(monthly_leads)
    natural_monthly_sales = nat_walk["sales"]

    # ─── Mode decision ──────────────────────────────────────────────────────
    MONTHLY_FLOOR = 4
    collapse = False
    if total_units > 0 and monthly_leads > 0:
        if total_units < 3:
            collapse = True
        elif natural_monthly_sales < MONTHLY_FLOOR:
            collapse = True
    # If we don't have both inputs, fall through to monthly presentation
    # (the legacy default) — no collapse, use nat_walk for the table.

    # ─── Build the "table walk" (what the table numbers represent) ─────────
    if collapse and total_units > 0:
        # Scale leads so derived sales == total_units (campaign-total walk).
        campaign_leads = round(total_units / SALE_PER_LEAD) if SALE_PER_LEAD else 0
        table_walk = _walk(campaign_leads)
    else:
        table_walk = nat_walk

    # Extract table-walk vars
    leads      = table_walk["leads"]
    appts      = table_walk["appts"]
    a_vis      = table_walk["a_vis"]
    nat        = table_walk["nat"]
    t_vis      = table_walk["t_vis"]
    offers     = table_walk["offers"]
    books      = table_walk["books"]
    sales      = table_walk["sales"]
    visit_to_sale_pct = (sales / t_vis * 100) if t_vis else 0
    overall_conv_pct  = (sales / leads * 100) if leads else 0
    leads_per_sale_t  = round(leads / sales) if sales else 0

    # ─── Header + sales_footer ──────────────────────────────────────────────
    header_title = "خطة المبيعات بالأرقام"
    if collapse:
        header_subtitle = f"قمع المبيعات للحملة — مشروع {_get(content,'project_name')}"
    else:
        header_subtitle = f"قمع المبيعات الشهري — مشروع {_get(content,'project_name')}"
    addHeader(s, header_title, header_subtitle)

    # ─── Right KPI sidebar ──────────────────────────────────────────────────
    _add_rect(s, 7.82, 1.22, 1.93, 4.15, BROWN)
    _add_rect(s, 7.82, 1.22, 1.93, 0.06, GOLD)
    _add_text(s, 7.85, 1.30, 1.85, 0.40, "مؤشرات التحويل",
              font_size=11, bold=True, color=GOLD, align="center", rtl=True)

    # months_to_sellout is ALWAYS derived from the natural monthly rate
    # (even when the table is collapsed). Sidebar is the honest schedule.
    if natural_monthly_sales > 0 and total_units > 0:
        months_val = max(1, round(total_units / natural_monthly_sales))
        months_str = f"{months_val}"
    else:
        # No inputs — fall back to content-dict value if provided
        months_str = str(_get(content, 'funnel_kpis', 'months_to_sellout',
                              default="—"))

    kpis = [
        (_get(content, 'funnel_kpis', 'overall_conversion'), "نسبة التحويل الإجمالية"),
        (_get(content, 'funnel_kpis', 'leads_per_sale'),     "مهتم لكل وحدة مباعة"),
        (_get(content, 'funnel_kpis', 'avg_unit_price'),     "متوسط سعر الوحدة"),
        (months_str,                                          "أشهر لبيع المتبقي من المشروع"),
    ]
    ys = [(1.78, 2.20), (2.63, 3.05), (3.48, 3.90), (4.33, 4.75)]
    for (num_y, lbl_y), (num, lbl) in zip(ys, kpis):
        _add_text(s, 7.85, num_y, 1.85, 0.42, num,
                  font_size=22, bold=True, color=GOLD, align="center", rtl=False)
        _add_text(s, 7.85, lbl_y, 1.85, 0.32, lbl,
                  font_size=10, color=SAND, align="center", rtl=True)

    # ─── LEFT: table ────────────────────────────────────────────────────────
    # Format numbers with thousands separator
    def _fmt(n): return f"{n:,}"
    def _pct2(p): return f"{p:.2f}%"  # for the fixed constants displayed as %
    def _pct1(p): return f"{p:.1f}%"  # for derived indicators

    # Row labels — one shifts in collapse mode (per user rule)
    sales_row_label = "إجمالي المبيعات" if collapse else "المبيعات"

    # Table rows: [label, value_as_string]
    rows = [
        ("المهتمين",                    _fmt(leads)),
        ("نسبة حجز موعد %",             _pct2(APPT_BOOKING_RATE * 100)),
        ("المواعيد",                    _fmt(appts)),
        ("نسبة الحضور %",               _pct2(APPT_ATTENDANCE * 100)),
        ("زيارات المواعيد",             _fmt(a_vis)),
        ("الزيارات الطبيعية",            _fmt(nat)),
        ("إجمالي الزيارات",              _fmt(t_vis)),
        ("المهتمين %",                   _pct2(INTERESTED_FROM_VIS * 100)),
        ("عروض الأسعار",                 _fmt(offers)),
        ("نسبة الحجز %",                 _pct2(BOOKING_RATE * 100)),
        ("الحجوزات",                    _fmt(books)),
        ("نسبة البيع من الحجوزات",      _pct2(SALE_FROM_BOOKING * 100)),
        (sales_row_label,               _fmt(sales)),
        ("نسبة الزيارات الى المبيعات",   _pct1(visit_to_sale_pct)),
        ("إجمال نسبة التحويل",           _pct1(overall_conv_pct)),
        ("عدد المهتمين للبيعة الوحدة",  _fmt(leads_per_sale_t)),
        ("إجمالي الوحدات المباعة",      _fmt(sales)),
    ]

    # Table layout — fits between left strip (x=0.18) and right sidebar (x=7.82)
    table_x = 0.35
    table_w = 7.40
    table_y = 1.22
    table_h = 4.15
    # Header row
    header_h = 0.35
    # Body rows fill the remainder
    n_rows = len(rows)
    row_h = (table_h - header_h) / n_rows

    # Column split: right half = label (physically right for RTL reading)
    #               left half = value (right-aligned to border)
    label_col_w = table_w * 0.60
    value_col_w = table_w * 0.40
    label_col_x = table_x + value_col_w  # label on the right (higher x)
    value_col_x = table_x                 # value on the left

    # Header
    _add_rect(s, table_x, table_y, table_w, header_h, COPPER)
    table_header = "الخطة البيعية للحملة" if collapse else "الخطة البيعية الشهرية"
    _add_text(s, table_x, table_y, table_w, header_h,
              table_header,
              font_size=13, bold=True, color=WHITE,
              align="center", rtl=True, anchor="middle")

    # Body rows with alternating fill
    for i, (label, value) in enumerate(rows):
        ry = table_y + header_h + i * row_h
        bg = RGBColor(0xFF, 0xF7, 0xF2) if i % 2 == 0 else WHITE
        _add_rect(s, table_x, ry, table_w, row_h, bg)
        # Row separator (thin sand line)
        _add_rect(s, table_x, ry + row_h - 0.005, table_w, 0.005, SAND)
        # Label (right-aligned, physically on the right column)
        _add_text(s, label_col_x + 0.10, ry, label_col_w - 0.20, row_h, label,
                  font_size=10.5, color=BROWN,
                  align="right", rtl=True, anchor="middle")
        # Value (centered in its narrower column)
        _add_text(s, value_col_x + 0.10, ry, value_col_w - 0.20, row_h, value,
                  font_size=11, bold=True, color=COPPER,
                  align="center", rtl=False, anchor="middle")

def slide_trapezoid(slide, x_top_left, y_top, w_top, x_bot_left, y_bot, w_bot, fill_rgb):
    """Draw a filled trapezoid by injecting a custom geometry <a:custGeom> into a shape.

    python-pptx's freeform builder is finicky with mid-page coordinates; direct XML
    is the reliable path. We start from a Rectangle shape and swap its preset
    geometry for a custom one defined by 4 points.
    """
    # Bounding box of the trapezoid
    bbox_left   = min(x_top_left, x_bot_left)
    bbox_right  = max(x_top_left + w_top, x_bot_left + w_bot)
    bbox_top    = y_top
    bbox_bottom = y_bot
    bbox_w = bbox_right - bbox_left
    bbox_h = bbox_bottom - bbox_top

    # Create a base rectangle we'll convert
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(bbox_left), Inches(bbox_top),
                                  Inches(bbox_w), Inches(bbox_h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    shp.line.fill.background()
    shp.shadow.inherit = False

    # Replace prstGeom with a custGeom defining the trapezoid.
    # Coordinate space: 0..100000 in both x and y (the default custGeom coord grid).
    sp = shp._element
    spPr = sp.find(qn("p:spPr"))
    for child in list(spPr):
        tag = child.tag.split("}")[-1]
        if tag == "prstGeom":
            spPr.remove(child)

    # Compute the 4 corners in 0..100000 space inside the bbox
    def nx(inches_val):
        return int((inches_val - bbox_left) / bbox_w * 100000)
    def ny(inches_val):
        return int((inches_val - bbox_top) / bbox_h * 100000)

    p_tl = (nx(x_top_left),          ny(y_top))
    p_tr = (nx(x_top_left + w_top),  ny(y_top))
    p_br = (nx(x_bot_left + w_bot),  ny(y_bot))
    p_bl = (nx(x_bot_left),          ny(y_bot))

    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    nsmap = {"a": a}
    custGeom = etree.SubElement(spPr, qn("a:custGeom"))
    etree.SubElement(custGeom, qn("a:avLst"))
    etree.SubElement(custGeom, qn("a:gdLst"))
    etree.SubElement(custGeom, qn("a:ahLst"))
    etree.SubElement(custGeom, qn("a:cxnLst"))
    rect = etree.SubElement(custGeom, qn("a:rect"),
                             l="0", t="0", r="100000", b="100000")
    pathLst = etree.SubElement(custGeom, qn("a:pathLst"))
    path = etree.SubElement(pathLst, qn("a:path"), w="100000", h="100000")

    moveTo = etree.SubElement(path, qn("a:moveTo"))
    etree.SubElement(moveTo, qn("a:pt"), x=str(p_tl[0]), y=str(p_tl[1]))
    for pt in (p_tr, p_br, p_bl):
        lnTo = etree.SubElement(path, qn("a:lnTo"))
        etree.SubElement(lnTo, qn("a:pt"), x=str(pt[0]), y=str(pt[1]))
    etree.SubElement(path, qn("a:close"))

    return shp

def slide15_why_wassel(prs, content):
    s = _blank_slide(prs)
    addStrip(s, variant="content")
    addHeader(s, "خطة المبيعات بالأرقام",
              f"قمع المبيعات الشهري — مشروع {_get(content,'project_name')}")
    addFooter(s, _get(content, 'sales_footer_right',
                      default=_get(content, 'footer_right')))

    # Right sidebar (same geometry as slide 14)
    _add_rect(s, 7.82, 1.22, 1.93, 4.15, BROWN)
    _add_rect(s, 7.82, 1.22, 1.93, 0.06, GOLD)
    _add_text(s, 7.85, 1.30, 1.85, 0.40, "لماذا وصل؟",
              font_size=11, bold=True, color=GOLD, align="center", rtl=True)

    # Sidebar bullets — NEVER write "Wassel CRM"; use "نظام وصل"
    sidebar_bullets = _get(content, 'why_wassel_sidebar', default=[
        "نظام وصل العقاري",
        "عملية مبيعات مؤتمة",
        "لوحة معلومات لمتابعة مباشرة",
        "إدارة الوحدات عن طريق نظام وصل",
    ])
    y = 1.85
    for b in sidebar_bullets:
        _add_text(s, 7.85, y, 1.85, 0.40,
                  f"◆ {b}",
                  font_size=10, color=SAND, align="right", rtl=True)
        y += 0.52

    # LEFT — reasons block on cream
    _add_text(s, 0.30, 1.30, 7.45, 0.48, "لماذا نظام وصل؟",
              font_size=16, bold=True, color=BROWN, align="right", rtl=True)

    reasons = _get(content, 'why_wassel_reasons', default=[
        ("تسجيل تلقائي",          "كل عميل يدخل النظام مباشرة — لا تضيع فرصة"),
        ("متابعة منظمة",          "إشعارات ومراحل تضمن عدم إهمال أي عميل"),
        ("لوحة معلومات فورية",     "مؤشرات أداء تظهر لحظة بلحظة للمالك"),
        ("إدارة الوحدات",          "حالة كل وحدة متاحة/محجوزة/مباعة محدّثة تلقائياً"),
        ("تقارير جاهزة",           "تقارير منتظمة تُرسل للمطوّر دون طلب"),
    ])
    y = 1.90
    for title, body in reasons:
        _add_rect(s, 0.30, y, 7.45, 0.60, WHITE)
        _add_rect(s, 7.67, y, 0.08, 0.60, COPPER)  # right-edge accent
        _add_text(s, 0.40, y + 0.05, 7.30, 0.24, f"◆  {title}",
                  font_size=12, bold=True, color=BROWN, align="right", rtl=True)
        _add_text(s, 0.60, y + 0.30, 7.10, 0.26, body,
                  font_size=10.5, color=CHARCOAL, align="right", rtl=True)
        y += 0.66

def slide16_closing(prs, content):
    s = _blank_slide(prs, bg_color=BROWN)
    _add_rect(s, 0.55, 0.00, 9.45, 0.04, GOLD)
    _add_rect(s, 0.55, 5.58, 9.45, 0.04, GOLD)
    # left copper strip (simplified version of divider strip)
    addStrip(s, variant="divider")
    # Overlay: the divider strip's dots visually work here too.

    _add_logo(s, 4.29, 0.24, 1.96, 1.96)
    _add_text(s, 0.70, 2.20, 9.10, 0.72, "شراكة تسويقية متكاملة",
              font_size=36, bold=True, color=CREAM, align="center", rtl=True)
    _add_text(s, 0.70, 2.90, 9.10, 0.42,
              "وصل العقارية —  النظام الذي يُحوّل الاهتمام إلى مبيعات",
              font_size=16, color=SAND, align="center", rtl=True)

    pillars = [
        (0.72, "إطلاق قوي",     "زخم من اليوم الأول"),
        (3.77, "تسويق مركّز",    "المحتوى والترويج يصلان للعميل الصحيح"),
        (6.82, "إغلاق منهجي",   "10 خطوات تُحوّل الاهتمام إلى حجوزات"),
    ]
    for x, title, body in pillars:
        _add_rect(s, x, 3.47, 2.82, 1.60, CREAM)
        _add_rect(s, x, 3.47, 2.82, 0.06, COPPER)
        _add_text(s, x + 0.10, 3.57, 2.62, 0.40, title,
                  font_size=14, bold=True, color=BROWN, align="center", rtl=True)
        _add_text(s, x + 0.10, 3.98, 2.62, 1.00, body,
                  font_size=12, color=CHARCOAL, align="center", rtl=True)

    _add_text(s, 0.70, 5.20, 9.10, 0.28, "wassel.re",
              font_size=13, color=GOLD, align="center", rtl=False,
              shape_hyperlink="https://wassel.re")

# ─── Entry point ──────────────────────────────────────────────────────────────

def build(content: dict, output_path: str):
    """Build a complete 15-slide Wassel deck from the content dict.

    `content` shape: see the module docstring SCHEMA.
    `output_path` should be an absolute path ending in .pptx.
    Returns the output_path.

    Slide sequence (15 slides):
       1 cover
       2 about Wassel
       3 divider — تحليل مربع المشروع
       4 market analysis
       5 competitor table
       6 project + opportunity
       7 divider — الخطة التسويقية
       8 opening event
       9 content & digital platforms
      10 marketing outcomes + derived targets
      11 divider — الخطة البيعية
      12 sales journey
      13 detailed sales steps
      14 sales plan by the numbers (formula-driven table)
      15 closing

    (The old slide 15 "why Wassel" was deleted — its content is implicit in
    the closing slide and KPI sidebar of slide 14. The `slide15_why_wassel`
    function is kept in the module for reference but not called.)
    """
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    slide1_cover(prs, content)
    slide2_about(prs, content)
    slide3_divider(prs, content)
    slide4_market(prs, content)
    slide5_competitors(prs, content)
    slide6_opportunity(prs, content)
    slide7_divider(prs, content)
    slide8_event(prs, content)
    slide9_content_platforms(prs, content)
    slide10_outcomes(prs, content)
    slide11_divider(prs, content)
    slide12_sales_snake(prs, content)
    slide13_sales_detailed(prs, content)
    slide14_funnel(prs, content)
    slide16_closing(prs, content)  # becomes slide 15 (closing) in the final deck

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    # Sample run with the reference مقام 17 data — regenerates a best-effort
    # clone of the reference deck for visual smoke-testing.
    sample = {
        "project_name": "مقام كورتيارد 17",
        "district":     "حي النرجس",
        "city":         "الرياض",
        "year":         "2026",
        "footer_right":           "وصل العقارية  |  مقام كورتيارد 17",
        "marketing_footer_right": "وصل العقارية  |  الخطة التسويقية",
        "sales_footer_right":     "وصل العقارية  |  الخطة البيعية",
        "market": {
            "kpis": [
                {"value": "457", "label": "إجمالي الوحدات في الحي"},
                {"value": "76%", "label": "متوسط نسبة البيع في الحي"},
                {"value": "7",   "label": "وحدات تُباع شهرياً في الحي"},
            ],
            "insight_lines": [
                "المشاريع المزوّدة بمرافق متكاملة حققت أعلى معدلات بيع في المربع — بين 80% و98%.",
                "هذا يُثبت أن المشتري يبحث عن أسلوب حياة لا مجرد وحدة سكنية.",
            ],
            "price_range_min": "830,000", "price_range_max": "1,500,000",
            "avg_price": "1,200,000",
        },
        "competitors": {
            "headers": ["المشروع", "المطوّر", "الوحدات", "نسبة البيع", "متوسط السعر", "ملاحظات"],
            "rows": [
                ["مقام 17", "—", "61", "~30%", "1.2M", "الإطلاق قبل 3 أشهر"],
            ],
            "footnote": "* انطلق المشروع قبل 3 أشهر فقط — 18 وحدة مباعة بمعدل 6 وحدات/شهر",
        },
        "project": {
            "location": "حي النرجس —  دقيقتين من طريق الملك سلمان",
            "area_range": "91 م² حتى 265 م²",
            "price_range": "918,400 الى 1,554,170 ريال",
            "unit_total": "61 وحدة — 4 عمائر \u200E(A/B/C/D)\u200E",
            "units_available": "~43 وحدة",
            "amenities": "نادي رياضي، لاونج، حضانة، غرفة سائق، مواقف",
            "warranties": "هيكل إنشائي 10 سنوات  |    عزل 15 سنة",
            "opportunity_bullets": [
                "الطلب مرتفع: الحي يسجل 7 وحدات مباعة شهرياً",
                "المرافق تُحرّك البيع — مقام 17 يملكها",
                "معدل مبيعات المشروع قريب من متوسط السوق",
                "43 وحدة متاحة بقيمة تجاوز 50 مليون ريال",
            ],
            "expected_revenue": "+50,000,000",
            "remaining_units": "43",
            "avg_price_millions": "1.2",
        },
        "marketing_targets": {
            # Only total_units_sold is still used; impressions/leads_target
            # are now DERIVED from it via fixed conversion rates.
            "total_units_sold": "43",
        },
        "sales_plan": {
            # Primary input for slide 14's formula-driven table. The whole
            # table derives from this monthly lead number.
            "leads_per_month": "1500",
        },
        "funnel_kpis": {
            "overall_conversion": "0.61%", "leads_per_sale": "166",
            "avg_unit_price": "1.3M",      "months_to_sellout": "3 - 6",
        },
    }
    out = os.path.expanduser("~/wassel_sample_output.pptx")
    build(sample, out)
    print(f"Built: {out}")
