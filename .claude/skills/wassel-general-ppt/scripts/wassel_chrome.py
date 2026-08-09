"""
Wassel General PPT — brand engine (lean).

This module gives you the Wassel BRAND in code form. It does NOT give you a
layout, a structure, a header treatment, a footer treatment, a card pattern,
or a slide size. Those are design decisions you make per deck.

What this module guarantees:
    - Brand palette (the only colors you should use)
    - Amiri font, applied across all OOXML font slots so Arabic actually renders
    - Arabic typography correctness (digit conversion, separators, em-dash
      bidi handling, line spacing) applied automatically to every string
    - Shape-level hyperlinks that don't get hijacked by the PowerPoint theme
    - Access to the white Wassel logo as an asset

What this module does NOT do:
    - Lock slide dimensions (you pick the size — 16:9, 4:3, A4, square, custom)
    - Provide cover/divider/header/footer/agenda/KPI/quote/etc. helpers
    - Suggest layout, spacing, card patterns, or section structures
    - Tell you when to use the logo, where to place text, or what to put on slides

Goal: every deck composed with this module should feel different. Same brand,
fresh design every time.


PUBLIC API
==========

Brand constants (RGBColor):
    COPPER   #B8734F
    SAND     #E8D9C0
    BROWN    #6B4226
    CREAM    #F8F5E9
    GOLD     #D9B57F
    CHARCOAL #3F3F3F
    WHITE    #FFFFFF
    FONT     "Amiri"

Common size presets (tuples of width, height in inches):
    SIZE_16_9, SIZE_4_3, SIZE_A4_LANDSCAPE, SIZE_A4_PORTRAIT, SIZE_SQUARE
    SIZE_INSTAGRAM_PORTRAIT, SIZE_INSTAGRAM_STORY
    (or just pass any (width, height) tuple — these are options, not requirements)

Unicode bidi marks:
    LRM, RLM, ALM, LRE, PDF, NBSP

Logo:
    LOGO_PATH  — absolute path to the white Wassel logo PNG

Functions:
    new_presentation(size=None, width=None, height=None)  — fresh Presentation
    blank_slide(prs, bg=CREAM)                            — slide w/ solid bg
    add_rect(slide, x, y, w, h, fill, line=None)          — filled rectangle
    add_text(slide, x, y, w, h, text, ...)                — Amiri text w/ all rules
    add_logo(slide, x, y, w, h)                           — drop the white logo
    add_shape_hyperlink(shape, url)                       — clickable shape
    clean_text(text)                                      — apply text fixes manually


TYPOGRAPHY RULES (auto-applied — you don't think about these)
=============================================================

Every string passed to add_text gets these fixes (set normalize=False to bypass):

  1. Western digits 0-9 -> Arabic-Indic ٠-٩.
     Reason: Arabic-Indic digits are strong-RTL; Western digits drift inside
     Arabic paragraphs and pull % / commas / periods to the wrong side.

  2. "." between digits -> "٫" (U+066B). "," between digits -> "٬" (U+066C).
     Reason: standalone "." and "," are bidi-neutral and drift in RTL text.
     The Arabic separators are strong-RTL and stay put. Standalone Arabic
     commas in sentences ("حي النرجس، الرياض") are left alone — no digit
     on either side, no replacement.

  3. Em-dash "—" in Arabic context -> wrapped with RLM marks.
     Reason: a bare em-dash is bidi-neutral and PowerPoint can absorb the
     surrounding spaces, rendering "حي النرجس —مدارس" instead of with
     proper spacing.

  4. Font set on latin + ea + cs OOXML slots.
     Reason: PowerPoint renders Arabic from the cs (complex script) slot.
     If empty, Arabic falls back to the theme font even though the font
     dropdown UI says Amiri. python-pptx's run.font.name only writes
     <a:latin> — we set ea and cs explicitly.

  5. Line spacing 1.5 when the block has >=8 words.
     Reason: long Arabic blocks with diacritics cramp without it.
     Override per-call with line_spacing=<float>.

  6. Hyperlinks via add_shape_hyperlink, NOT run.hyperlink.address.
     Reason: run-level hyperlinks trigger PowerPoint's theme to override
     run color with blue + add an underline. Shape-level links attach the
     URL to the shape's <p:cNvPr> instead — click works, theme stays out.


WORDING RULES (NOT auto-enforced — check manually)
==================================================

  - Always "نادي" — never "نادٍ".
  - Always "نظام وصل" — never "Wassel CRM" / "CRM وصل".
  - Latin building codes inside Arabic paragraphs wrapped in LRM marks:
    f"4 عمائر {LRM}(A/B/C/D){LRM}".
"""

import os
import re as _re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# --- Brand constants ---------------------------------------------------------
COPPER   = RGBColor(0xB8, 0x73, 0x4F)
SAND     = RGBColor(0xE8, 0xD9, 0xC0)
BROWN    = RGBColor(0x6B, 0x42, 0x26)
CREAM    = RGBColor(0xF8, 0xF5, 0xE9)
GOLD     = RGBColor(0xD9, 0xB5, 0x7F)
CHARCOAL = RGBColor(0x3F, 0x3F, 0x3F)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Amiri"


# --- Size presets (options, not requirements) -------------------------------
# Pick one for new_presentation(size=...) or pass any custom (w, h) tuple.
SIZE_16_9              = (10.0, 5.625)
SIZE_4_3               = (10.0, 7.5)
SIZE_A4_LANDSCAPE      = (11.69, 8.27)
SIZE_A4_PORTRAIT       = (8.27, 11.69)
SIZE_SQUARE            = (10.0, 10.0)
SIZE_INSTAGRAM_PORTRAIT = (8.0, 10.0)   # 4:5 portrait
SIZE_INSTAGRAM_STORY   = (5.625, 10.0)  # 9:16 story


# --- Unicode bidi marks ------------------------------------------------------
LRM  = "\u200E"   # Left-to-right mark
RLM  = "\u200F"   # Right-to-left mark
ALM  = "\u061C"   # Arabic letter mark
LRE  = "\u202A"   # LTR embedding open
PDF  = "\u202C"   # Pop directional format
NBSP = "\u00A0"


# --- Logo asset --------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.abspath(os.path.join(HERE, "..", "assets", "wassel_logo_white.png"))


# --- Slide creation ----------------------------------------------------------

def new_presentation(size=None, width=None, height=None):
    """Fresh Presentation. You must pick the canvas size.

    Pass either a `size` tuple (e.g. SIZE_16_9, SIZE_A4_PORTRAIT, or any custom
    `(w, h)` tuple in inches), OR explicit `width=` and `height=` kwargs.

    No default — the choice of canvas is part of the design. Pick what fits
    the deck's purpose: presentation-on-screen, printed handout, social post,
    portrait mobile, square Instagram, etc.
    """
    if size is not None:
        w, h = size
    elif width is not None and height is not None:
        w, h = width, height
    else:
        raise ValueError(
            "new_presentation requires a size. Pass size=SIZE_16_9 (or any "
            "other preset), or pass a custom (w, h) tuple, or explicit "
            "width=... height=... kwargs in inches."
        )
    prs = Presentation()
    prs.slide_width = Inches(w)
    prs.slide_height = Inches(h)
    return prs


def blank_slide(prs, bg=CREAM):
    """New slide with a solid background covering the whole canvas.

    bg: any RGBColor — brand constant or custom. Background dimensions are
    read from the presentation, so this works on any canvas size.

    Set bg=None to skip the background paint (transparent / theme default).
    """
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    if bg is not None:
        w_in = prs.slide_width / Emu(1) / 914400.0
        h_in = prs.slide_height / Emu(1) / 914400.0
        add_rect(slide, 0, 0, w_in, h_in, bg)
    return slide


# --- Drawing primitives ------------------------------------------------------

def add_rect(slide, x, y, w, h, fill, line=None):
    """Filled rectangle. Coordinates and dimensions in inches.

    fill: RGBColor.
    line: RGBColor for outline, or None for no outline (default).
    """
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_logo(slide, x, y, w, h):
    """Place the white Wassel logo at (x, y) with size (w, h). All in inches.
    No-ops silently if the logo asset is missing."""
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))


# --- Text engine (typography rules live here) --------------------------------

def _set_rtl(paragraph):
    """Set paragraph direction to RTL. python-pptx doesn't expose this."""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")


def _set_font_all_scripts(run, font_name):
    """Set Latin, East-Asian, AND Complex-Script font on a run.

    PowerPoint reads Arabic from the cs slot. python-pptx's run.font.name
    only writes <a:latin>. If cs is empty, Arabic falls back to the theme
    default — even though the UI font dropdown says the right name.
    """
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    for local in ("ea", "cs"):
        for existing in rPr.findall(qn(f"a:{local}")):
            rPr.remove(existing)
    latin = rPr.find(qn("a:latin"))
    insert_at = list(rPr).index(latin) + 1 if latin is not None else len(rPr)
    for local in ("ea", "cs"):
        elem = etree.Element(qn(f"a:{local}"))
        elem.set("typeface", font_name)
        rPr.insert(insert_at, elem)
        insert_at += 1


def _normalize_spacing(text):
    """Clean spacing around separators. For em/en-dashes in Arabic context,
    wrap in RLM marks to keep surrounding spaces from collapsing."""
    if not isinstance(text, str) or not text:
        return text
    text = _re.sub(r" {2,}", " ", text)
    has_arabic = any("\u0600" <= ch <= "\u06FF" for ch in text)
    for sep in ("—", "–"):
        pat_sep = _re.escape(sep)
        if has_arabic:
            text = _re.sub(rf"\s*{pat_sep}\s*", f" {RLM}{sep}{RLM} ", text)
        else:
            text = _re.sub(rf"\s*{pat_sep}\s*", f" {sep} ", text)
    text = _re.sub(r"\s*\|\s*", " | ", text)
    text = _re.sub(r" {2,}", " ", text)
    return text.strip()


_WESTERN_TO_ARABIC_INDIC = str.maketrans("0123456789", "٠١٢٣٤٥٦٧٨٩")


def _convert_digits(text):
    """Convert Western digits to Arabic-Indic, and replace decimal/thousands
    separators between digits with their Arabic equivalents (٫, ٬)."""
    if not isinstance(text, str) or not text:
        return text
    text = text.translate(_WESTERN_TO_ARABIC_INDIC)
    text = _re.sub(r"(?<=[٠-٩])\.(?=[٠-٩])", "٫", text)
    text = _re.sub(r"(?<=[٠-٩]),(?=[٠-٩])", "٬", text)
    return text


def clean_text(text):
    """Apply all Wassel typography fixes to a string. Useful when building
    strings for table cells or other contexts where add_text doesn't apply."""
    return _convert_digits(_normalize_spacing(text))


def _strip_hyperlink_style(run, color):
    """Defeat PowerPoint's theme override of hyperlink runs (blue + underline).
    Used internally when add_text is called with a run-level hyperlink.
    Prefer add_shape_hyperlink, which sidesteps this entirely.
    """
    rPr = run._r.get_or_add_rPr()
    hlink = rPr.find(qn("a:hlinkClick"))
    if hlink is None:
        return
    rPr.set("u", "none")
    for existing in rPr.findall(qn("a:uFill")):
        rPr.remove(existing)
    uFill = etree.SubElement(rPr, qn("a:uFill"))
    etree.SubElement(uFill, qn("a:noFill"))
    for existing in hlink.findall(qn("a:solidFill")):
        hlink.remove(existing)
    solidFill = etree.SubElement(hlink, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", str(color))


def add_text(slide, x, y, w, h, text, *,
             font_size=12, bold=False, color=CHARCOAL,
             align="right", rtl=True, font_name=FONT,
             anchor="middle", hyperlink=None, shape_hyperlink=None,
             normalize=True, line_spacing=None):
    """Drop-in textbox with full Arabic typography handling.

    text:           str OR list[str] — each string becomes one paragraph.
    font_size:      int (pt).
    bold:           bool.
    color:          RGBColor — pick from brand constants.
    align:          "right" | "center" | "left".
    rtl:            bool — True for Arabic, False for pure Latin.
    anchor:         "top" | "middle" | "bottom" — vertical anchor in the box.
    hyperlink:      URL for a run-level hyperlink. AVOID — PowerPoint will
                    theme-override the color/underline. Use shape_hyperlink.
    shape_hyperlink: URL for a shape-level hyperlink. Click target is the
                    whole textbox; theme styling does NOT apply. Preferred.
    normalize:      bool — apply digit/separator/dash cleanups (default True).
    line_spacing:   float override. If None, auto-rule applies: 1.5 when block
                    has >=8 words total, else default.

    Returns the textbox shape.
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map[anchor]

    lines = text if isinstance(text, list) else [text]

    if line_spacing is None:
        def _strip_marks(s):
            return "".join(ch for ch in s if ch not in (LRM, RLM, ALM, LRE, PDF))
        total_words = sum(len(_strip_marks(ln).split()) for ln in lines)
        line_spacing = 1.5 if total_words >= 8 else None

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    for i, ln in enumerate(lines):
        if normalize:
            ln = clean_text(ln)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[align]
        if rtl:
            _set_rtl(p)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        _set_font_all_scripts(run, font_name)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        if hyperlink:
            run.hyperlink.address = hyperlink
            _strip_hyperlink_style(run, color)

    if shape_hyperlink:
        add_shape_hyperlink(tb, shape_hyperlink)

    return tb


def add_shape_hyperlink(shape, url):
    """Attach a hyperlink to the SHAPE (not a text run). Click works; the
    PowerPoint theme does NOT override the text color/underline."""
    slide_part = shape.part
    rId = slide_part.relate_to(
        url,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    sp = shape._element
    nvSpPr = sp.find(qn("p:nvSpPr"))
    if nvSpPr is None:
        return
    cNvPr = nvSpPr.find(qn("p:cNvPr"))
    if cNvPr is None:
        return
    for existing in cNvPr.findall(qn("a:hlinkClick")):
        cNvPr.remove(existing)
    hlink = etree.SubElement(cNvPr, qn("a:hlinkClick"))
    hlink.set(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id",
        rId,
    )


__all__ = [
    # Brand
    "COPPER", "SAND", "BROWN", "CREAM", "GOLD", "CHARCOAL", "WHITE", "FONT",
    # Sizes (options)
    "SIZE_16_9", "SIZE_4_3", "SIZE_A4_LANDSCAPE", "SIZE_A4_PORTRAIT",
    "SIZE_SQUARE", "SIZE_INSTAGRAM_PORTRAIT", "SIZE_INSTAGRAM_STORY",
    # Bidi marks
    "LRM", "RLM", "ALM", "LRE", "PDF", "NBSP",
    # Logo
    "LOGO_PATH",
    # Functions
    "new_presentation", "blank_slide",
    "add_rect", "add_text", "add_logo", "add_shape_hyperlink", "clean_text",
]
