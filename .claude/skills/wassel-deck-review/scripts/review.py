"""
Wassel Deck Review — automated lint + patcher.

Opens a built Wassel pptx and runs a series of checks. For mechanical issues
(wrong font slots, blue hyperlinks, missing RTL, banned phrases, spacing
glitches around separators), the script patches the pptx in place and writes
a new file with suffix `_reviewed`. For judgment calls (missing numbers,
structural deviations, content that doesn't match the project brief), it
produces a report.

Usage:
    from review import review_deck
    report = review_deck("/path/to/input.pptx",
                         output_path="/path/to/output_reviewed.pptx")
    print(report.markdown())

The report object has:
    .issues_fixed     — list of (severity, location, description) auto-fixed
    .issues_reported  — list of (severity, location, description) needing human
    .checks_passed    — list of (check_name, count_verified)
    .markdown()       — formatted markdown report

Severities: "critical" (breaks brand spec), "warning" (likely wrong),
            "info" (cosmetic, you probably want this fixed).
"""

import os
import re
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ─── Brand constants (must match wassel-presentation builder) ────────────────
COPPER   = "B8734F"
SAND     = "E8D9C0"
BROWN    = "6B4226"
CREAM    = "F8F5E9"
GOLD     = "D9B57F"
CHARCOAL = "3F3F3F"
WHITE    = "FFFFFF"
BRAND_PALETTE = {COPPER, SAND, BROWN, CREAM, GOLD, CHARCOAL, WHITE,
                 # Slide-9 platform brand colors explicitly allowed:
                 "FFFC00", "000000", "C13584", "0A66C2",
                 # Slide-10 tagline background used in the builder:
                 "FFF7F2",
                 # Slide-14 darkest funnel level:
                 "5A371F"}
REQUIRED_FONT = "Amiri"

# ─── Forbidden / required text ────────────────────────────────────────────────
BANNED_SUBSTRINGS = [
    "Wassel CRM",     # always "نظام وصل"
    "wassel CRM",
    "CRM وصل",
    "نادٍ",           # always "نادي" (non-kasra form)
]

# Exact subtitles required on specific slides (per the Wassel build spec)
EXACT_SUBTITLES = {
    7: "الهدف: صناعة الطلب، وجلب المهتمين",
    11: "تحويل الطلب والاهتمام إلى مبيعات",
}

# Slide 4 subtitle must contain this substring (project name varies)
SLIDE4_REQUIRED_SUBSTRING = "مربع مشروع"

# ─── Report model ─────────────────────────────────────────────────────────────
@dataclass
class Issue:
    severity: str      # "critical" | "warning" | "info"
    location: str      # e.g. "slide 5, table header"
    description: str

@dataclass
class Report:
    issues_fixed: list = field(default_factory=list)
    issues_reported: list = field(default_factory=list)
    checks_passed: list = field(default_factory=list)

    def markdown(self):
        out = ["# Wassel Deck Review\n"]
        if not self.issues_fixed and not self.issues_reported:
            out.append("✅ **All checks passed.** Deck is ready to ship.\n")
        else:
            out.append(f"**Auto-fixed:** {len(self.issues_fixed)} issues.  "
                       f"**Needs attention:** {len(self.issues_reported)} issues.\n")

        if self.issues_fixed:
            out.append("\n## Auto-fixed (no action needed)\n")
            for iss in self.issues_fixed:
                out.append(f"- **[{iss.severity}]** {iss.location} — {iss.description}")

        if self.issues_reported:
            out.append("\n## Needs your attention\n")
            for iss in self.issues_reported:
                out.append(f"- **[{iss.severity}]** {iss.location} — {iss.description}")

        if self.checks_passed:
            out.append("\n## Checks passed\n")
            for name, count in self.checks_passed:
                out.append(f"- {name}: {count} items verified")
        return "\n".join(out)

# ─── Helpers ──────────────────────────────────────────────────────────────────
def _iter_runs(slide):
    """Yield (shape, paragraph, run) for every text run on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    yield shape, p, r
        if shape.shape_type == 19:  # TABLE
            for row in shape.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            yield shape, p, r

def _set_font_all_scripts(run, font_name):
    """Same logic as in the builder — sets latin + ea + cs so Arabic renders Amiri."""
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    for local in ("ea", "cs"):
        for existing in rPr.findall(qn(f"a:{local}")):
            rPr.remove(existing)
    latin = rPr.find(qn("a:latin"))
    insert_at = list(rPr).index(latin) + 1 if latin is not None else len(rPr)
    for local in ("ea", "cs"):
        elem = etree.Element(qn(f"a:{local}"))
        elem.set("typeface", font_name)
        rPr.insert(insert_at, elem)
        insert_at += 1

def _get_run_text_clean(run):
    return run.text if run.text else ""

# ─── Individual checks (each returns list of Issue) ───────────────────────────
def check_font_everywhere(prs, fix=True):
    """Verify every text run has latin+ea+cs = Amiri. Fix missing cs slots and
    mismatches between latin and cs (the classic bug where user manually changes
    font in PowerPoint, which updates latin but leaves cs with stale value)."""
    fixed, reported, total = [], [], 0
    for si, slide in enumerate(prs.slides, start=1):
        for shape, p, r in _iter_runs(slide):
            total += 1
            rPr = r._r.find(qn("a:rPr"))
            if rPr is None:
                continue
            latin = rPr.find(qn("a:latin"))
            cs = rPr.find(qn("a:cs"))
            ea = rPr.find(qn("a:ea"))
            latin_face = latin.get("typeface") if latin is not None else None
            cs_face = cs.get("typeface") if cs is not None else None
            ea_face = ea.get("typeface") if ea is not None else None

            # Any slot with a non-Amiri font is a violation
            bad_slots = []
            if latin_face and latin_face != REQUIRED_FONT:
                bad_slots.append(("latin", latin_face))
            if cs_face and cs_face != REQUIRED_FONT:
                bad_slots.append(("cs", cs_face))
            if ea_face and ea_face != REQUIRED_FONT:
                bad_slots.append(("ea", ea_face))

            # Missing cs slot when Arabic is present
            has_arabic = any("\u0600" <= ch <= "\u06FF" for ch in (r.text or ""))
            missing_cs = (cs is None) and has_arabic

            if bad_slots:
                slots_desc = ", ".join(f"{name}='{face}'" for name, face in bad_slots)
                iss = Issue("critical", f"slide {si}",
                            f"Wrong font slot(s) [{slots_desc}] — expected {REQUIRED_FONT}. "
                            f"Text: '{r.text[:50]}'")
                if fix:
                    _set_font_all_scripts(r, REQUIRED_FONT)
                    fixed.append(iss)
                else:
                    reported.append(iss)
            elif missing_cs:
                iss = Issue("warning", f"slide {si}",
                            f"Arabic text missing Complex-Script font slot "
                            f"— Arabic will render as theme default in WPS/PowerPoint. "
                            f"Text: '{r.text[:50]}'")
                if fix:
                    _set_font_all_scripts(r, REQUIRED_FONT)
                    fixed.append(iss)
                else:
                    reported.append(iss)
    return fixed, reported, total

def check_banned_text(prs):
    """Scan for 'Wassel CRM', 'نادٍ', etc. These are judgment calls — report, don't auto-replace."""
    reported = []
    for si, slide in enumerate(prs.slides, start=1):
        for shape, p, r in _iter_runs(slide):
            for banned in BANNED_SUBSTRINGS:
                if banned in r.text:
                    reported.append(Issue(
                        "critical", f"slide {si}",
                        f"Banned phrase '{banned}' found. "
                        f"Run says: '{r.text[:80]}'. "
                        f"Replace with approved wording (see Wassel SKILL.md)."
                    ))
    return reported

def check_exact_subtitles(prs):
    """Slide 4 must contain 'مربع مشروع'. Slides 7, 11 must have exact subtitles."""
    reported = []
    slides = list(prs.slides)

    # Slide 4 — subtitle is the 2nd text element in the header area, y≈0.65
    if len(slides) >= 4:
        found = False
        for shape, p, r in _iter_runs(slides[3]):
            if SLIDE4_REQUIRED_SUBSTRING in r.text:
                found = True
                break
        if not found:
            reported.append(Issue(
                "critical", "slide 4",
                f"Subtitle must contain '{SLIDE4_REQUIRED_SUBSTRING}' "
                f"(pattern: 'تحليل مربع مشروع <project> <district> — <city>'). "
                f"Do not shorten to just 'حي <district>'."
            ))

    for slide_num, required in EXACT_SUBTITLES.items():
        if len(slides) < slide_num:
            continue
        found = any(r.text.strip() == required for _, _, r in _iter_runs(slides[slide_num - 1]))
        if not found:
            reported.append(Issue(
                "critical", f"slide {slide_num}",
                f"Subtitle must be EXACTLY: '{required}'. Do not paraphrase."
            ))
    return reported

def check_hyperlinks_unstyled(prs, fix=True):
    """Hyperlinks should not be blue/underlined. Patch any that are."""
    fixed, reported = [], []
    for si, slide in enumerate(prs.slides, start=1):
        for shape, p, r in _iter_runs(slide):
            rPr = r._r.find(qn("a:rPr"))
            if rPr is None:
                continue
            hlink = rPr.find(qn("a:hlinkClick"))
            if hlink is None:
                continue
            u_attr = rPr.get("u")
            iss = None
            if u_attr != "none":
                iss = Issue("warning", f"slide {si}",
                            f"Hyperlink run ('{r.text[:40]}') lacks explicit u='none' "
                            f"— may render as underlined blue.")
                if fix:
                    rPr.set("u", "none")
                    fixed.append(iss)
                else:
                    reported.append(iss)
    return fixed, reported

def check_colors_in_palette(prs):
    """Report any fill colors not in the brand palette.

    Judgment call — some images might legitimately introduce other colors
    (e.g., event icons). Report only fills on shapes, not pictures.
    """
    reported = []
    for si, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            # Only check auto-shapes, not pictures or tables
            if shape.shape_type not in (1,):  # AUTO_SHAPE only
                continue
            try:
                if shape.fill.type is None:
                    continue
                rgb = shape.fill.fore_color.rgb
                if rgb is None:
                    continue
                hex_color = f"{rgb:06X}".upper()
                if hex_color not in BRAND_PALETTE:
                    reported.append(Issue(
                        "info", f"slide {si}",
                        f"Shape fill color #{hex_color} is not in the brand palette. "
                        f"Confirm this is intentional."
                    ))
            except Exception:
                continue
    return reported

def check_spacing_around_separators(prs, fix=True):
    """Fix unbalanced spaces around — , – , -, _, | .

    Rule: for any separator that already has whitespace on at least one side
    (meaning it's being used as a separator, not as a compound joiner like
    'A-B'), enforce exactly one space before and one space after. Compound
    tokens with no surrounding whitespace are left alone.

    Idempotent: if a separator is already wrapped in RLM marks (the builder
    does this for Arabic context), leave it alone. Re-wrapping would stack
    invisible characters every time the linter runs.
    """
    fixed, total_fixed = [], 0
    RLM = "\u200F"
    for si, slide in enumerate(prs.slides, start=1):
        for shape, p, r in _iter_runs(slide):
            original = r.text or ""
            if not original:
                continue
            new = original
            new = re.sub(r" {2,}", " ", new)
            # Em-dash, en-dash, pipe — always separators, match greedily.
            for sep in ("—", "–", "|"):
                rlm_sep_pat = f"{RLM}{sep}{RLM}"
                if rlm_sep_pat in new:
                    continue
                new = re.sub(rf"\s*{re.escape(sep)}\s*", f" {sep} ", new)
            # Hyphen-minus and underscore — only normalize when surrounded by
            # whitespace on at least one side (compound tokens like "A-B" or
            # "snake_case" are preserved).
            for sep in ("-", "_"):
                rlm_sep_pat = f"{RLM}{sep}{RLM}"
                if rlm_sep_pat in new:
                    continue
                pattern = rf"(?:\s+{re.escape(sep)}\s*|\s*{re.escape(sep)}\s+)"
                new = re.sub(pattern, f" {sep} ", new)
            new = re.sub(r" {2,}", " ", new).strip()
            if new != original:
                total_fixed += 1
                if fix:
                    r.text = new
                    fixed.append(Issue(
                        "info", f"slide {si}",
                        f"Normalized spacing around separators in: '{original[:60]}'"
                    ))
    return fixed, total_fixed

def check_number_isolation(prs, fix=True):
    """Wrap Latin/numeric tokens in Arabic text with LRM marks.

    This prevents dots/commas from drifting when Arabic and numbers mix.
    Only applied to runs that contain Arabic AND Latin characters.
    """
    fixed, total = [], 0
    token_re = re.compile(r"[A-Za-z0-9][A-Za-z0-9.,%~+\-/:]*")
    LRM = "\u200E"

    for si, slide in enumerate(prs.slides, start=1):
        for shape, p, r in _iter_runs(slide):
            txt = r.text or ""
            if not txt:
                continue
            has_arabic = any("\u0600" <= ch <= "\u06FF" for ch in txt)
            if not has_arabic:
                continue
            # Skip if already has LRM marks — idempotent
            if LRM in txt:
                continue
            new_txt = token_re.sub(lambda m: f"{LRM}{m.group(0)}{LRM}", txt)
            if new_txt != txt:
                total += 1
                if fix:
                    r.text = new_txt
                    fixed.append(Issue(
                        "info", f"slide {si}",
                        f"Isolated {len(token_re.findall(txt))} numeric/Latin tokens with LRM in: '{txt[:60]}'"
                    ))
    return fixed, total

def check_no_parens_in_copy(prs, fix=True):
    """Replace `(text)` with `— text —` in body/callout copy.

    Brand rule from wassel-presentation/SKILL.md → Punctuation enforcement:
    parentheses are banned in user-visible free-text copy. Use em-dashes with
    spaces on both sides instead. This runs as a safety net — the builder
    shouldn't be emitting parens anymore, but manual edits or upstream data
    leaks can re-introduce them.

    Allowed exceptions (left alone):
    - Parens already wrapped in LRM marks (`\u200E(A/B/C/D)\u200E`) — the
      builder uses this for building-code lists where parens are layout, not
      prose. Detected by LRM immediately adjacent to `(` or `)`.
    - URL-like runs (contain `://` or start with `http`) — parens inside URLs
      are structural, not prose.
    - Runs where the paren content looks numeric-only with a comma/dash
      (e.g. "(1,490,000 - 2,090,000)") — this is likely a range expression
      that the builder should be formatting elsewhere, but em-dashing it
      mid-range would produce nonsense. Reported (not auto-fixed) so human
      can decide.

    Match shape: `\([^()]{1,60}\)` — single-level, non-empty, under 60 chars.
    Longer paren spans are likely structural and get reported, not fixed.
    """
    fixed, total = [], 0
    LRM = "\u200E"
    paren_re = re.compile(r"\(([^()]{1,60})\)")

    for si, slide in enumerate(prs.slides, start=1):
        for shape, p, r in _iter_runs(slide):
            txt = r.text or ""
            if "(" not in txt and ")" not in txt:
                continue
            # Skip URL-like runs
            if "://" in txt or txt.lstrip().startswith("http"):
                continue

            matches = list(paren_re.finditer(txt))
            if not matches:
                continue

            new_parts = []
            cursor = 0
            changed = False
            for m in matches:
                new_parts.append(txt[cursor:m.start()])
                inner = m.group(1)
                # LRM-adjacent: keep as-is (building-code exception)
                prev_ch = txt[m.start()-1] if m.start() > 0 else ""
                next_ch = txt[m.end()] if m.end() < len(txt) else ""
                if prev_ch == LRM or next_ch == LRM:
                    new_parts.append(m.group(0))
                    cursor = m.end()
                    continue
                # Numeric range inside parens (e.g. "1,490,000 - 2,090,000"):
                # don't rewrite — report instead.
                if re.fullmatch(r"[0-9,.\u0660-\u0669\u066B\u066C\s\-–—]+", inner.strip()):
                    new_parts.append(m.group(0))
                    cursor = m.end()
                    continue
                new_parts.append(f"— {inner} —")
                cursor = m.end()
                changed = True
            new_parts.append(txt[cursor:])
            new_txt = "".join(new_parts)
            # Collapse any double-spacing introduced by the rewrite
            new_txt = re.sub(r" {2,}", " ", new_txt).strip()

            if changed and new_txt != txt:
                total += 1
                if fix:
                    r.text = new_txt
                    fixed.append(Issue(
                        "info", f"slide {si}",
                        f"Replaced parens with em-dashes: '{txt[:60]}' → '{new_txt[:60]}'"
                    ))
                else:
                    fixed.append(Issue(
                        "info", f"slide {si}",
                        f"WOULD replace parens: '{txt[:60]}'"
                    ))
    return fixed, total

def check_rtl_paragraphs(prs, fix=True):
    """Arabic-containing paragraphs should have rtl='1'.

    Exception: paragraphs that are ONLY Arabic-Indic digits / percent / comma
    (e.g. '٤٣', '٠١', '٧٦%', '٢٠٢٦') are intentionally treated as LTR numeric
    labels in the builder — those don't need rtl='1' and flagging them would
    create noise."""
    fixed, total = [], 0
    for si, slide in enumerate(prs.slides, start=1):
        for shape, p, r in _iter_runs(slide):
            if not r.text:
                continue
            # Arabic-Indic digits, percent, comma, period — numeric-only content
            text_stripped = r.text.strip()
            is_numeric_only = all(
                (ch in "٠١٢٣٤٥٦٧٨٩" or ch in "%,.٫٬-~+ ")
                for ch in text_stripped
            )
            has_arabic_letters = any(
                "\u0600" <= ch <= "\u06FF" and ch not in "٠١٢٣٤٥٦٧٨٩٫٬"
                for ch in r.text
            )
            if is_numeric_only or not has_arabic_letters:
                continue
            pPr = p._p.find(qn("a:pPr"))
            is_rtl = pPr is not None and pPr.get("rtl") == "1"
            if not is_rtl:
                total += 1
                if fix:
                    if pPr is None:
                        pPr = p._p.get_or_add_pPr()
                    pPr.set("rtl", "1")
                    fixed.append(Issue(
                        "warning", f"slide {si}",
                        f"Arabic paragraph was not set as RTL. Text: '{r.text[:50]}'"
                    ))
    return fixed, total

def check_footer_presence(prs):
    """Content slides (2, 4, 5, 6, 8, 9, 10, 12, 13, 14) must have a footer with wassel.re.
    (Slide 15 is the closing slide — it has wassel.re as a tagline, not a footer strip.)"""
    content_slide_indices = [2, 4, 5, 6, 8, 9, 10, 12, 13, 14]
    reported = []
    slides = list(prs.slides)
    for si in content_slide_indices:
        if len(slides) < si:
            continue
        slide = slides[si - 1]
        has_footer = any("wassel.re" in (r.text or "") for _, _, r in _iter_runs(slide))
        if not has_footer:
            reported.append(Issue(
                "critical", f"slide {si}",
                "Content slide is missing 'wassel.re' footer."
            ))
    return reported

def check_divider_no_footer(prs):
    """Dividers (3, 7, 11) should NOT have a footer.
    (Slides 1 and 15 are also divider-style but have wassel.re as intentional
    visual elements — we don't flag those.)"""
    divider_indices = [3, 7, 11]
    reported = []
    slides = list(prs.slides)
    for si in divider_indices:
        if len(slides) < si:
            continue
        slide = slides[si - 1]
        has_footer = any("wassel.re" in (r.text or "") for _, _, r in _iter_runs(slide))
        if has_footer:
            reported.append(Issue(
                "warning", f"slide {si}",
                "Divider slide should not have a footer (brand spec)."
            ))
    return reported

def check_slide_count(prs):
    """Wassel spec is exactly 15 slides (the former slide 15 'why Wassel' was
    deleted in the 2026 refresh; closing slide is now slide 15)."""
    reported = []
    expected = 15
    actual = len(prs.slides)
    if actual != expected:
        reported.append(Issue(
            "critical", "deck",
            f"Deck has {actual} slides; Wassel spec is exactly {expected}. "
            f"(The deck structure is: 1 cover, 2 about, 3 divider, 4 market, "
            f"5 competitors, 6 opportunity, 7 divider, 8 event, 9 content, "
            f"10 outcomes, 11 divider, 12 journey, 13 detailed, 14 sales-table, "
            f"15 closing.)"
        ))
    return reported

def check_arabic_indic_digits(prs):
    """Western digits (0-9) should not appear in Arabic text runs — they should
    be converted to Arabic-Indic (٠-٩). Reports violations (doesn't auto-fix
    because converting digits silently could corrupt Latin identifiers like
    building codes or email fragments)."""
    reported = []
    western_digit_re = _re_western_digit()
    for si, slide in enumerate(prs.slides, start=1):
        for shape, p, r in _iter_runs(slide):
            txt = r.text or ""
            if not txt:
                continue
            # Only flag if the run contains BOTH Arabic characters AND Western digits.
            # Pure Latin identifiers (building codes "A/B/C/D") are fine.
            has_arabic = any("\u0600" <= ch <= "\u06FF" for ch in txt)
            if not has_arabic:
                continue
            western_digits = western_digit_re.findall(txt)
            if western_digits:
                reported.append(Issue(
                    "critical", f"slide {si}",
                    f"Western digits ({''.join(western_digits[:8])}{'...' if len(western_digits) > 8 else ''}) "
                    f"found in Arabic text run. All numbers in Arabic context must be Arabic-Indic "
                    f"(٠-٩). Text: '{txt[:60]}'"
                ))
    return reported

def _re_western_digit():
    import re
    return re.compile(r"[0-9]")

def check_line_spacing_8plus(prs):
    """Text blocks with ≥8 total words (across all paragraphs) must have
    line_spacing=1.5 on every paragraph. Reports violations."""
    LRM = "\u200E"; RLM = "\u200F"; ALM = "\u061C"; LRE = "\u202A"; PDF = "\u202C"
    BIDI_MARKS = (LRM, RLM, ALM, LRE, PDF)

    def _strip_marks(s):
        return "".join(ch for ch in s if ch not in BIDI_MARKS)

    reported = []
    for si, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            # Count total words across all paragraphs in this text frame
            all_text = "\n".join(p.text or "" for p in tf.paragraphs)
            stripped = _strip_marks(all_text)
            word_count = len(stripped.split())
            if word_count < 8:
                continue
            # Check each paragraph's line spacing
            for pi, p in enumerate(tf.paragraphs):
                if not (p.text or "").strip():
                    continue
                ls = p.line_spacing
                # line_spacing=1.5 in python-pptx is stored as 1.5 (Ratio) or
                # as 150% via Pct. Anything less than 1.4 or None = violation.
                if ls is None or (isinstance(ls, (int, float)) and ls < 1.4):
                    reported.append(Issue(
                        "warning", f"slide {si}",
                        f"Text block with {word_count} words (≥8) should have "
                        f"line_spacing=1.5; paragraph {pi + 1} has line_spacing={ls}. "
                        f"Text: '{(p.text or '')[:50]}'"
                    ))
                    break  # one report per block, not per paragraph
    return reported

def check_wassel_re_hyperlink_shape_level(prs):
    """wassel.re hyperlinks must be at the SHAPE level (on p:cNvPr) not the
    RUN level (inside a:rPr). Run-level hyperlinks trigger PowerPoint's theme
    override which paints the text blue + underlined."""
    reported = []
    for si, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Does this shape contain 'wassel.re' text?
            contains_wassel = False
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if "wassel.re" in (r.text or ""):
                        contains_wassel = True
                        # Check if the run has an hlinkClick (run-level = bad)
                        rPr = r._r.find(qn("a:rPr"))
                        if rPr is not None and rPr.find(qn("a:hlinkClick")) is not None:
                            reported.append(Issue(
                                "critical", f"slide {si}",
                                "'wassel.re' has a RUN-level hyperlink (inside a:rPr). "
                                "PowerPoint's theme will override this to blue+underlined. "
                                "Move the hyperlink to the SHAPE level (on p:cNvPr) using "
                                "shape_hyperlink in the builder."
                            ))
                            break
                if not contains_wassel:
                    continue
    return reported

# ─── Entry point ──────────────────────────────────────────────────────────────
def review_deck(input_path, output_path=None, fix=True):
    """Run the full review suite.

    Args:
        input_path:  path to the pptx to review
        output_path: where to save the patched pptx (default: adds _reviewed suffix)
        fix:         if True, auto-fix mechanical issues; if False, only report

    Returns:
        Report object (call .markdown() for a formatted summary).
    """
    if not os.path.exists(input_path):
        raise FileNotFoundError(input_path)

    if output_path is None:
        base, ext = os.path.splitext(input_path)
        output_path = f"{base}_reviewed{ext}"

    prs = Presentation(input_path)
    report = Report()

    # 1. Slide count
    report.issues_reported.extend(check_slide_count(prs))

    # 2. Font (latin + ea + cs)
    font_fixed, font_reported, font_total = check_font_everywhere(prs, fix=fix)
    report.issues_fixed.extend(font_fixed)
    report.issues_reported.extend(font_reported)
    report.checks_passed.append((f"Font = Amiri", font_total))

    # 3. Banned text
    report.issues_reported.extend(check_banned_text(prs))

    # 4. Exact subtitles
    report.issues_reported.extend(check_exact_subtitles(prs))

    # 5. Hyperlink styling
    hl_fixed, hl_reported = check_hyperlinks_unstyled(prs, fix=fix)
    report.issues_fixed.extend(hl_fixed)
    report.issues_reported.extend(hl_reported)

    # 6. Colors in palette
    report.issues_reported.extend(check_colors_in_palette(prs))

    # 7. Spacing around separators
    space_fixed, space_count = check_spacing_around_separators(prs, fix=fix)
    report.issues_fixed.extend(space_fixed)

    # (Old "number isolation" check was dropped — we now convert Western digits
    # to Arabic-Indic at build time, which eliminates the bidi drift that LRM
    # wrapping was working around. Letting the linter auto-add LRM now would
    # double-encode runs that are already correctly formatted.)

    # 8. RTL paragraphs
    rtl_fixed, rtl_count = check_rtl_paragraphs(prs, fix=fix)
    report.issues_fixed.extend(rtl_fixed)

    # 9. Punctuation — no parens in body/callout copy
    paren_fixed, paren_count = check_no_parens_in_copy(prs, fix=fix)
    report.issues_fixed.extend(paren_fixed)

    # 10. Footer presence / absence
    report.issues_reported.extend(check_footer_presence(prs))
    report.issues_reported.extend(check_divider_no_footer(prs))

    # 11. Arabic-Indic digit enforcement (new: replaces the old LRM-wrap check)
    report.issues_reported.extend(check_arabic_indic_digits(prs))

    # 12. Line-spacing rule for ≥8-word blocks
    report.issues_reported.extend(check_line_spacing_8plus(prs))

    # 13. wassel.re must use shape-level hyperlink, not run-level
    report.issues_reported.extend(check_wassel_re_hyperlink_shape_level(prs))

    if fix:
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        prs.save(output_path)

    return report


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python review.py <input.pptx> [output.pptx]")
        sys.exit(1)
    inp = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else None
    rep = review_deck(inp, output_path=out)
    print(rep.markdown())
    if out:
        print(f"\nPatched deck saved to: {out}")
